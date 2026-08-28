"""
Automatização Slides Semanais — Banco Bari
Interface web para geração dos slides de diretoria.

Deploy: Streamlit Cloud (https://streamlit.io/cloud)
"""

import streamlit as st
import os, io, calendar, tempfile, json, base64
from collections import defaultdict
from datetime import datetime, date, timedelta
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.path import Path
import matplotlib.patches as mpatches
import openpyxl
from pptx import Presentation
from pptx.util import Inches
import requests

# ══════════════════════════════════════════════════════════════
# PAGE CONFIG & CSS
# ══════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Slides Semanais — Banco Bari",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed",
)

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;500;600;700;800;900&display=swap');

    /* ═══ BARI BRAND TOKENS ═══ */
    :root {
        --bari-blue: #4A90E2;
        --bari-blue-dark: #2563EB;
        --bari-blue-light: #EBF2FC;
        --bari-navy: #0A1628;
        --bari-navy-2: #0D1B2A;
        --bari-navy-3: #142236;
        --bari-gray-50: #F5F7FA;
        --bari-gray-100: #E8ECF2;
        --bari-gray-200: #D0D6E0;
        --bari-gray-400: #8E99A8;
        --bari-gray-600: #556070;
        --bari-gray-800: #2D3142;
        --bari-white: #FFFFFF;
        --bari-orange: #F59E0B;
        --bari-red: #E53E3E;
        --bari-green-ok: #10B981;
    }

    /* ═══ GLOBAL ═══ */
    .stApp {
        font-family: 'Outfit', 'Segoe UI', system-ui, sans-serif !important;
        background: var(--bari-gray-50) !important;
    }
    .block-container { max-width: 940px !important; padding-top: 0 !important; }
    header[data-testid="stHeader"] { display: none !important; }
    #MainMenu { display: none !important; }
    footer { display: none !important; }
    div[data-testid="stDecoration"] { display: none !important; }

    /* ═══ HEADER ═══ */
    .bari-header {
        background: var(--bari-navy-2);
        padding: 0;
        margin: -1rem -1rem 32px -1rem;
        color: white;
        position: relative;
        overflow: hidden;
    }
    .bari-header-inner {
        display: flex;
        align-items: center;
        justify-content: space-between;
        padding: 22px 36px;
        position: relative;
        z-index: 2;
    }
    .bari-header::after {
        content: '';
        position: absolute;
        bottom: 0; left: 0; right: 0;
        height: 3px;
        background: linear-gradient(90deg, var(--bari-blue) 0%, var(--bari-blue-dark) 50%, transparent 100%);
    }
    /* Decorative dots pattern */
    .bari-header::before {
        content: '';
        position: absolute;
        top: 0; right: 0;
        width: 300px; height: 100%;
        background: radial-gradient(circle at 2px 2px, rgba(74,144,226,0.08) 1px, transparent 0);
        background-size: 20px 20px;
    }
    .bari-logo {
        font-size: 32px;
        font-weight: 900;
        letter-spacing: -1.5px;
        color: white;
        line-height: 1;
    }
    .bari-logo span {
        color: var(--bari-blue);
    }
    .bari-header-title {
        font-size: 14px;
        font-weight: 500;
        color: rgba(255,255,255,0.5);
        margin-top: 2px;
        letter-spacing: 0.3px;
    }
    .bari-header-right {
        display: flex;
        align-items: center;
        gap: 8px;
        background: rgba(255,255,255,0.06);
        padding: 8px 16px;
        border-radius: 8px;
        border: 1px solid rgba(255,255,255,0.08);
    }
    .bari-header-right span {
        font-size: 12px;
        color: rgba(255,255,255,0.45);
        font-weight: 500;
    }
    .bari-header-dot {
        width: 7px; height: 7px;
        border-radius: 50%;
        background: var(--bari-blue);
        animation: pulse-dot 2s ease-in-out infinite;
    }
    @keyframes pulse-dot {
        0%, 100% { opacity: 1; }
        50% { opacity: 0.4; }
    }

    /* ═══ SECTION HEADERS ═══ */
    .step-header {
        display: flex; align-items: center; gap: 12px;
        margin-bottom: 16px; margin-top: 8px;
    }
    .step-num {
        width: 30px; height: 30px; border-radius: 8px;
        background: var(--bari-blue);
        color: white;
        display: flex; align-items: center; justify-content: center;
        font-size: 14px; font-weight: 700;
        flex-shrink: 0;
    }
    .step-num-inactive {
        background: var(--bari-gray-200) !important;
        color: var(--bari-gray-400) !important;
    }
    .step-title {
        font-weight: 700; font-size: 16px; color: var(--bari-navy);
    }
    .step-sub {
        font-size: 12px; color: var(--bari-gray-400); margin-left: 4px;
    }

    /* ═══ CARDS BASE ═══ */
    .bari-card {
        background: var(--bari-white);
        border-radius: 12px;
        border: 1px solid var(--bari-gray-100);
        padding: 20px 22px;
        transition: all 0.2s ease;
    }
    .bari-card:hover {
        border-color: var(--bari-gray-200);
        box-shadow: 0 2px 12px rgba(0,0,0,0.04);
    }

    /* ═══ FILE UPLOAD AREA ═══ */
    .stFileUploader > div > div { padding: 6px !important; }
    div[data-testid="stFileUploaderDropzone"] {
        padding: 10px !important;
        border-radius: 10px !important;
        border-color: var(--bari-gray-200) !important;
        background: var(--bari-gray-50) !important;
    }
    div[data-testid="stFileUploaderDropzone"]:hover {
        border-color: var(--bari-blue) !important;
        background: var(--bari-blue-light) !important;
    }
    div[data-testid="stFileUploaderDropzone"] > div > span {
        font-size: 12px !important;
        color: var(--bari-gray-400) !important;
    }
    div[data-testid="stFileUploaderDropzone"] button {
        font-size: 12px !important;
        padding: 5px 14px !important;
        border-radius: 8px !important;
        background: var(--bari-blue) !important;
        color: white !important;
        border: none !important;
    }

    /* ═══ BADGES ═══ */
    .badge-ok {
        background: var(--bari-blue-light); color: var(--bari-blue-dark);
        padding: 2px 10px; border-radius: 6px;
        font-size: 10px; font-weight: 700;
        letter-spacing: 0.5px;
    }
    .badge-req {
        background: #FFF3E0; color: var(--bari-orange);
        padding: 2px 10px; border-radius: 6px;
        font-size: 10px; font-weight: 700;
        letter-spacing: 0.5px;
    }

    /* ═══ DATE CARDS ═══ */
    .date-card {
        background: var(--bari-white);
        border-radius: 10px;
        padding: 14px 16px;
        border: 1px solid var(--bari-gray-100);
        border-top: 3px solid;
    }
    .date-card-green  { border-top-color: var(--bari-blue); }
    .date-card-blue   { border-top-color: #3B82F6; }
    .date-card-purple { border-top-color: #8B5CF6; }
    .date-label {
        font-size: 10px; color: var(--bari-gray-400); font-weight: 600;
        text-transform: uppercase; letter-spacing: 0.8px;
        margin-bottom: 6px;
    }
    .date-value { font-size: 20px; font-weight: 800; }
    .date-green  .date-value { color: var(--bari-blue-dark); }
    .date-blue   .date-value { color: #2563EB; }
    .date-purple .date-value { color: #7C3AED; }

    /* ═══ SUMMARY ═══ */
    .summary-box {
        background: var(--bari-white); border-radius: 10px;
        padding: 14px 18px; border: 1px solid var(--bari-gray-100);
        font-size: 13px; color: var(--bari-gray-600);
        line-height: 1.5; margin-bottom: 14px;
        border-left: 3px solid var(--bari-blue);
    }

    /* ═══ LOG PANEL ═══ */
    .log-panel {
        background: var(--bari-navy-2);
        border-radius: 10px;
        padding: 16px 18px;
        max-height: 340px;
        overflow-y: auto;
        font-family: 'JetBrains Mono', 'Fira Code', 'Consolas', monospace;
        font-size: 11.5px;
        line-height: 1.7;
        border: 1px solid rgba(74,144,226,0.15);
    }
    .log-normal  { color: #cbd5e1; }
    .log-success { color: #4ade80; }
    .log-warning { color: #fbbf24; }
    .log-error   { color: #f87171; }
    .log-info    { color: #60a5fa; }

    /* ═══ SUCCESS BANNER ═══ */
    .success-banner {
        background: linear-gradient(135deg, var(--bari-blue-light) 0%, #DBEAFE 100%);
        border-radius: 12px;
        padding: 28px 32px;
        border: 2px solid var(--bari-blue);
        text-align: center;
    }

    /* ═══ NOTE BOX ═══ */
    .note-box {
        background: #FFF8F0;
        border-radius: 10px;
        padding: 14px 18px;
        border: 1px solid #FFDDB5;
        border-left: 3px solid var(--bari-orange);
        font-size: 13px;
        color: #7C4A1E;
        line-height: 1.5;
    }

    /* ═══ BUTTONS ═══ */
    .stButton > button[kind="primary"],
    .stDownloadButton > button {
        background: var(--bari-blue) !important;
        color: white !important;
        border-radius: 10px !important;
        font-weight: 700 !important;
        letter-spacing: 0.3px !important;
        padding: 14px 24px !important;
        font-size: 15px !important;
        border: none !important;
        transition: all 0.2s ease !important;
    }
    .stButton > button[kind="primary"]:hover,
    .stDownloadButton > button:hover {
        background: var(--bari-blue-dark) !important;
        box-shadow: 0 4px 16px rgba(74,144,226,0.3) !important;
    }

    /* ═══ TOGGLE ═══ */
    .stToggle label span { font-size: 13px !important; }

    /* ═══ EXPANDER ═══ */
    .streamlit-expanderHeader {
        font-size: 14px !important;
        font-weight: 600 !important;
        color: var(--bari-gray-600) !important;
    }

    /* ═══ DIVIDER ═══ */
    .soft-divider {
        height: 1px;
        background: linear-gradient(90deg, var(--bari-gray-100) 0%, transparent 100%);
        margin: 28px 0;
    }

    /* ═══ LABEL STYLING ═══ */
    .upload-label {
        font-weight: 700;
        font-size: 13px;
        color: var(--bari-navy);
        margin-bottom: 2px;
    }
    .upload-label-req {
        display: inline-block;
        background: #FFF3E0;
        color: var(--bari-orange);
        font-size: 9px;
        font-weight: 700;
        padding: 1px 6px;
        border-radius: 4px;
        margin-left: 6px;
        vertical-align: middle;
        letter-spacing: 0.5px;
    }
</style>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════
# CONSTANTES
# ══════════════════════════════════════════════════════════════

FASES_ORDEM = [
    ("Novo",                                   "Data de criação"),
    ("Tentativa de contato",                   "Data Etapa Tentativa de contato"),
    ("Trabalhando/negociação",                 "Data Etapa Trabalhando/negociação"),
    ("Aguardando documentação",                "Data Etapa Aguardando documentação"),
    ("Pré-análise",                            "Data Etapa Pré-análise"),
    ("Análise de crédito",                     "Data Etapa Análise de Crédito"),
    ("Crédito aprovado",                       "Data Etapa Crédito aprovado"),
    ("Análise jurídica / Avaliação do imóvel", "Data Etapa Análise Jurídica"),
    ("Emissão do contrato",                    "Data Etapa Emissão de contrato"),
    ("Assinatura",                             "Data Etapa Assinatura"),
]

FASES_NOMES = [f for f, _ in FASES_ORDEM]

CORES = {
    "Novo":                                   "#4472C4",
    "Tentativa de contato":                   "#1F3864",
    "Trabalhando/negociação":                 "#70D4CE",
    "Aguardando documentação":                "#2E8B6A",
    "Pré-análise":                            "#D4BE6A",
    "Análise de crédito":                     "#E8922A",
    "Crédito aprovado":                       "#C0392B",
    "Análise jurídica / Avaliação do imóvel": "#C97A1A",
    "Emissão do contrato":                    "#27AE60",
    "Assinatura":                             "#1A5276",
}

CANAIS = {
    "B2C":           ["B2C"],
    "GP":            ["GP"],
    "PC":            ["PC"],
    "Relacionamento":["Relacionamento"],
    "Comercial":     ["B2C", "GP", "PC"],
    "Todos":         ["B2C", "GP", "PC", "Relacionamento"],
}

FASES_PC = [f for f in FASES_NOMES if f not in (
    "Novo", "Tentativa de contato", "Trabalhando/negociação", "Aguardando documentação"
)]

SLIDES_FUNIL = [
    (9,  "B2C",       "volume",    "mensal"),
    (10, "B2C",       "volume",    "semanal"),
    (11, "B2C",       "propostas", "mensal"),
    (12, "B2C",       "propostas", "semanal"),
    (14, "GP",        "volume",    "mensal"),
    (15, "GP",        "volume",    "semanal"),
    (16, "GP",        "propostas", "mensal"),
    (17, "GP",        "propostas", "semanal"),
    (19, "PC",        "volume",    "mensal"),
    (20, "PC",        "volume",    "semanal"),
    (21, "PC",        "propostas", "mensal"),
    (22, "PC",        "propostas", "semanal"),
    (24, "Comercial", "volume",    "mensal"),
    (25, "Comercial", "volume",    "semanal"),
    (26, "Comercial", "propostas", "mensal"),
    (27, "Comercial", "propostas", "semanal"),
]

POS_ESQ     = (1.26, 1.21, 3.24, 4.39)
POS_DIR     = (4.70, 1.21, 3.24, 4.39)
POS_LEGENDA = (7.40, 3.40, 2.50, 2.20)

SLIDES_DASH = [
    (13, "B2C"),
    (18, "GP"),
    (23, "PC"),
    (28, "Relacionamento"),
    (29, "Todos"),
]

POS_DASH = (1.20, 0.74, 8.03, 4.85)

EXCLUIR_TIMES = {"GP": ["FRANQ"], "Todos": ["FRANQ"]}

DIST_MTD = {m: {1: 0.20, 2: 0.45, 3: 0.70, 4: 0.90, 5: 1.00} for m in range(1, 13)}

PLAN_STAGE_MAP = {
    "lead": "Lead", "wl": "Workable Lead", "workable lead": "Workable Lead",
    "novo": "Novo consultor", "trabalhando": "Trabalhando",
    "aguardando documentação": "Documentação", "aguardando documentacao": "Documentação",
    "pré-análise": "Pré-análise", "pre-analise": "Pré-análise",
    "análise de crédito": "Análise crédito", "analise de credito": "Análise crédito",
    "crédito aprovado": "Crédito aprovado", "credito aprovado": "Crédito aprovado",
    "jurídico/imóvel": "Jurídica", "jurídico/ imóvel": "Jurídica",
    "juridico/imovel": "Jurídica", "juridico/ imovel": "Jurídica",
    "emissão": "Emissão", "emissao": "Emissão",
    "assinatura": "Assinatura", "efetivado": "Novos contratos",
}

PLAN_TABS = {"B2C": "B2C", "GP": "GP", "PC": "PC", "Rel": "Relacionamento", "Total CGI": "Todos"}

PLAN_MESES = {
    "janeiro": 1, "fevereiro": 2, "março": 3, "marco": 3,
    "abril": 4, "maio": 5, "junho": 6, "julho": 7, "agosto": 8,
    "setembro": 9, "outubro": 10, "novembro": 11, "dezembro": 12,
}

MESES_PT = {
    1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
    5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
    9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro",
}

FASES_DASH = {
    "B2C": ["Lead","Workable Lead","Novo consultor","Trabalhando","Documentação","Pré-análise","Análise crédito","Crédito aprovado","Jurídica","Emissão","Assinatura","Novos contratos"],
    "GP": ["Workable Lead","Novo consultor","Trabalhando","Documentação","Pré-análise","Análise crédito","Crédito aprovado","Jurídica","Emissão","Assinatura","Novos contratos"],
    "PC": ["Novo consultor","Pré-análise","Análise crédito","Crédito aprovado","Jurídica","Emissão","Assinatura","Novos contratos"],
    "Relacionamento": ["Novo consultor","Pré-análise","Análise crédito","Crédito aprovado","Jurídica","Emissão","Assinatura","Novos contratos"],
    "Todos": ["Lead","Workable Lead","Novo consultor","Trabalhando","Documentação","Pré-análise","Análise crédito","Crédito aprovado","Jurídica","Emissão","Assinatura","Novos contratos"],
}

CONVERSOES_DASH = {
    "B2C": [("Trabalhando → Efetivado","Trabalhando","Novos contratos"),("Novo → Pré-Análise","Novo consultor","Pré-análise"),("Crédito Aprov. → Efetivado","Crédito aprovado","Novos contratos")],
    "GP": [("Trabalhando → Efetivado","Trabalhando","Novos contratos"),("Novo → Pré-Análise","Novo consultor","Pré-análise"),("Crédito Aprov. → Efetivado","Crédito aprovado","Novos contratos")],
    "PC": [("Pré-análise → Efetivado","Pré-análise","Novos contratos"),("Novo → Pré-Análise","Novo consultor","Pré-análise"),("Crédito Aprov. → Efetivado","Crédito aprovado","Novos contratos")],
    "Relacionamento": [("Pré-análise → Efetivado","Pré-análise","Novos contratos"),("Novo → Pré-Análise","Novo consultor","Pré-análise"),("Crédito Aprov. → Efetivado","Crédito aprovado","Novos contratos")],
    "Todos": [("Novo consultor → Efetivado","Novo consultor","Novos contratos"),("Novo → Pré-Análise","Novo consultor","Pré-análise"),("Crédito Aprov. → Efetivado","Crédito aprovado","Novos contratos")],
}

DASH_FASE_TO_OPP = {
    "Novo consultor": "Novo", "Trabalhando": "Trabalhando/negociação",
    "Documentação": "Aguardando documentação", "Pré-análise": "Pré-análise",
    "Análise crédito": "Análise de crédito", "Crédito aprovado": "Crédito aprovado",
    "Jurídica": "Análise jurídica / Avaliação do imóvel",
    "Emissão": "Emissão do contrato", "Assinatura": "Assinatura",
}

DASH_DOT_CORES = {
    "Lead": "#3b82f6", "Workable Lead": "#2563eb", "Novo consultor": "#7c3aed",
    "Trabalhando": "#0891b2", "Documentação": "#0e7490", "Pré-análise": "#059669",
    "Análise crédito": "#047857", "Crédito aprovado": "#065f46", "Jurídica": "#d97706",
    "Emissão": "#b45309", "Assinatura": "#92400e", "Novos contratos": "#166534",
}

STAGE_COL_DF = {
    'Novo consultor': 'Data de criação', 'Trabalhando': 'Data Etapa Trabalhando/negociação',
    'Documentação': 'Data Etapa Aguardando documentação', 'Pré-análise': 'Data Etapa Pré-análise',
    'Análise crédito': 'Data Etapa Análise de Crédito', 'Crédito aprovado': 'Data Etapa Crédito aprovado',
    'Jurídica': 'Data Etapa Análise Jurídica', 'Emissão': 'Data Etapa Emissão de contrato',
    'Assinatura': 'Data Etapa Assinatura',
}

DATE_COLS_OPPS = [
    'Data de criação', 'Data Etapa Trabalhando/negociação',
    'Data Etapa Aguardando documentação', 'Data Etapa Pré-análise',
    'Data Etapa Análise de Crédito', 'Data Etapa Crédito aprovado',
    'Data Etapa Análise Jurídica', 'Data Etapa Emissão de contrato',
    'Data Etapa Assinatura', 'Data de fechamento',
    'Data da última mudança de fase',
]

# ══════════════════════════════════════════════════════════════
# FUNÇÕES DE PROCESSAMENTO
# ══════════════════════════════════════════════════════════════

def parse_data(val):
    if val is None or val == '':
        return None
    if isinstance(val, datetime):
        return val.date()
    if isinstance(val, date):
        return val
    s = str(val).strip()
    for fmt in ('%d/%m/%Y %H:%M', '%d/%m/%Y', '%Y-%m-%d'):
        try:
            return datetime.strptime(s, fmt).date()
        except ValueError:
            continue
    return None


def carregar_base(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
    ws = wb.active
    hdrs = [c.value for c in next(ws.iter_rows(min_row=1, max_row=1))]
    col = {h: i for i, h in enumerate(hdrs) if h}
    rows = []
    for r in ws.iter_rows(min_row=2, values_only=True):
        fase = r[col['Fase']]
        time_ = r[col['Time']]
        if not fase or not time_:
            continue
        valor = float(r[col['Valor do Derivado']] or 0)
        data_fech = parse_data(r[col['Data de fechamento']])
        datas = {}
        for fase_nome, col_nome in FASES_ORDEM:
            if col_nome in col:
                datas[fase_nome] = parse_data(r[col[col_nome]])
        rows.append({'time': time_, 'fase': fase, 'valor': valor, 'datas': datas, 'data_fechamento': data_fech})
    wb.close()
    return rows


def carregar_planejamento(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    metas = {}
    for tab_nome, canal in PLAN_TABS.items():
        if tab_nome not in wb.sheetnames:
            continue
        ws = wb[tab_nome]
        mes_cell = None
        for row in ws.iter_rows():
            for cell in row:
                if str(cell.value or "").strip().lower() in ("mês", "mes"):
                    mes_cell = cell
                    break
            if mes_cell:
                break
        if not mes_cell:
            continue
        stage_col = mes_cell.column
        month_row = mes_cell.row
        col_mes_map = {}
        for cell in ws[month_row]:
            if cell.column <= stage_col:
                continue
            val = str(cell.value or "").strip().lower()
            if val in PLAN_MESES:
                col_mes_map[cell.column] = PLAN_MESES[val]
        metas[canal] = {}
        for row in ws.iter_rows(min_row=month_row + 1):
            nome_cell = ws.cell(row=row[0].row, column=stage_col)
            raw = str(nome_cell.value or "").strip().lower()
            fase_nome = PLAN_STAGE_MAP.get(raw)
            if not fase_nome:
                continue
            for col_idx, mes_num in col_mes_map.items():
                cell = ws.cell(row=nome_cell.row, column=col_idx)
                if cell.value is None:
                    continue
                try:
                    val_int = int(float(cell.value))
                except (ValueError, TypeError):
                    continue
                metas[canal].setdefault(mes_num, {})[fase_nome] = val_int
    wb.close()
    return metas


def retrato_funil(rows, canal, data_ref):
    times = CANAIS[canal]
    count = defaultdict(int)
    volume = defaultdict(float)
    for row in rows:
        if row['time'] not in times:
            continue
        if row['fase'] in ('Fechado ganho', 'Fechado perdido'):
            if row['data_fechamento'] and row['data_fechamento'] <= data_ref:
                continue
        fase_na_data = None
        for fase_nome, _ in FASES_ORDEM:
            dt = row['datas'].get(fase_nome)
            if dt is not None and dt <= data_ref:
                fase_na_data = fase_nome
        if fase_na_data:
            count[fase_na_data] += 1
            volume[fase_na_data] += row['valor']
    return count, volume


def fmt_valor(v):
    if v >= 1_000_000: return f"R${v/1_000_000:.0f}m"
    if v >= 1_000: return f"R${v/1_000:.0f}K"
    return f"R${v:,.0f}"


def gerar_funil_png(count, volume, canal, tipo, data_label):
    # Usa todas as fases, mas filtra as que têm valor zero
    fases_base = FASES_NOMES
    if tipo == "volume":
        fases = [f for f in fases_base if volume.get(f, 0) > 0]
        valores = [volume.get(f, 0) for f in fases]
    else:
        fases = [f for f in fases_base if count.get(f, 0) > 0]
        valores = [count.get(f, 0) for f in fases]
    total = sum(valores)
    n = len(fases)
    fig, ax = plt.subplots(figsize=(3.2, 4.2))
    fig.patch.set_facecolor('white')
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    ax.text(0.50, 0.99, data_label, ha='center', va='top', fontsize=9, fontweight='bold', color='#222222')
    titulo = f"Soma de Valor do Derivado: R${total/1e6:.0f}m" if tipo == "volume" else f"Contagem de registros: {int(total)}"
    ax.text(0.50, 0.91, titulo, ha='center', va='top', fontsize=7.0, color='#555555')
    if n == 0 or total == 0:
        ax.text(0.5, 0.5, 'Sem dados', ha='center', va='center', fontsize=10, color='#999')
        buf = io.BytesIO(); plt.savefig(buf, format='png', dpi=180, bbox_inches='tight', facecolor='white'); plt.close(); buf.seek(0); return buf.read()
    ft, fb = 0.84, 0.04; fh = ft - fb; top_hw, bot_hw = 0.45, 0.17; cx = 0.50
    props = [v / total for v in valores]; cum = 0.0
    for fase, val, prop in zip(fases, valores, props):
        yt = ft - fh * cum; yb = ft - fh * (cum + prop)
        hw_top = top_hw * (1 - cum) + bot_hw * cum
        hw_bot = top_hw * (1 - (cum+prop)) + bot_hw * (cum+prop)
        verts = [(cx-hw_top,yt),(cx+hw_top,yt),(cx+hw_bot,yb),(cx-hw_bot,yb),(cx-hw_top,yt)]
        codes = [Path.MOVETO, Path.LINETO, Path.LINETO, Path.LINETO, Path.CLOSEPOLY]
        patch = mpatches.PathPatch(Path(verts, codes), facecolor=CORES[fase], edgecolor='none')
        ax.add_patch(patch)
        pct = prop * 100; ym = (yt + yb) / 2; wm = hw_top + hw_bot; sh = yt - yb
        txt = f"{fmt_valor(val)} ({pct:.2f}%)" if tipo == "volume" else f"{int(val)} ({pct:.2f}%)"
        fs = 7.5 if wm > 0.50 else 6.5 if wm > 0.38 else 5.5
        if val > 0 and sh > 0.03:
            ax.text(cx, ym, txt, ha='center', va='center', fontsize=fs, fontweight='bold', color='black')
        cum += prop
    plt.tight_layout(pad=0.2)
    buf = io.BytesIO(); plt.savefig(buf, format='png', dpi=180, bbox_inches='tight', facecolor='white'); plt.close(); buf.seek(0); return buf.read()


def gerar_legenda_png(fases_list=None):
    fases = fases_list if fases_list else FASES_NOMES
    fig, ax = plt.subplots(figsize=(2.4, 2.6)); fig.patch.set_facecolor('white')
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    ax.text(0.95, 0.97, "Fase", ha='right', va='top', fontsize=9, color='#444444', fontweight='bold')
    n = len(fases); spacing = min(0.087, 0.85 / max(n, 1))
    for i, fase in enumerate(fases):
        y = 0.89 - i * spacing
        ax.add_patch(plt.Circle((0.90, y), 0.022, color=CORES[fase], transform=ax.transData, clip_on=False))
        nome = fase.replace("Análise jurídica / Avaliação do imóvel", "Análise jurídica / Av. imóvel")
        ax.text(0.83, y, nome, ha='right', va='center', fontsize=8, color='#444444')
    plt.tight_layout(pad=0.2)
    buf = io.BytesIO(); plt.savefig(buf, format='png', dpi=180, bbox_inches='tight', facecolor='white'); plt.close(); buf.seek(0); return buf.read()


def carregar_opps_df(file_bytes):
    df = pd.read_excel(io.BytesIO(file_bytes), dtype=str)
    for col in DATE_COLS_OPPS:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], dayfirst=True, errors='coerce')
    return df


def carregar_leads_df(file_bytes):
    if file_bytes is None: return None
    df = pd.read_excel(io.BytesIO(file_bytes), dtype=str)
    if 'Data de criação' in df.columns:
        df['Data de criação'] = pd.to_datetime(df['Data de criação'], dayfirst=True, errors='coerce')
    return df


def calcular_ref_periodos(df):
    ref_cols = [c for c in ['Data de criação','Data Etapa Trabalhando/negociação','Data da última mudança de fase'] if c in df.columns]
    REF = df[ref_cols].max(numeric_only=False).max()
    ano, mes, dia = REF.year, REF.month, REF.day
    re_start = pd.Timestamp(ano, mes, 1); re_end = REF.normalize()
    sw_start = pd.Timestamp(ano, mes, 1); sw_end = (REF - pd.Timedelta(days=7)).normalize()
    prev_mes = mes - 1 if mes > 1 else 12; prev_ano = ano if mes > 1 else ano - 1
    ultimo_dia_prev = calendar.monthrange(prev_ano, prev_mes)[1]
    mp_start = pd.Timestamp(prev_ano, prev_mes, 1); mp_end = pd.Timestamp(prev_ano, prev_mes, min(dia, ultimo_dia_prev))
    return REF, re_start, re_end, mp_start, mp_end, sw_start, sw_end


def filter_opps_df(df, canal):
    if 'Time' not in df.columns: return df
    if canal == 'Todos': return df
    if canal == 'Relacionamento': return df[df['Time'].isin(['Relacionamento', 'Rel'])]
    if canal in ('B2C', 'PC', 'GP'): return df[df['Time'] == canal]
    return df


def count_stage_df(df, stage, start, end):
    id_col = 'ID da oportunidade' if 'ID da oportunidade' in df.columns else None
    if stage == 'Novos contratos':
        if 'Data de fechamento' not in df.columns: return 0
        mask = (df['Data de fechamento'] >= start) & (df['Data de fechamento'] <= end) & (df['Fase'] == 'Fechado ganho')
        return int(df[mask][id_col].nunique()) if id_col else int(mask.sum())
    col = STAGE_COL_DF.get(stage)
    if not col or col not in df.columns: return 0
    mask = (df[col] >= start) & (df[col] <= end)
    return int(df[mask][id_col].nunique()) if id_col else int(mask.sum())


def count_leads_df(df_leads, canal, start, end):
    if df_leads is None or len(df_leads) == 0: return {'Lead': 0, 'Workable Lead': 0}
    df_c = df_leads
    if canal not in ('Todos',) and 'Canal' in df_leads.columns:
        if canal in ('B2C', 'GP'): df_c = df_leads[df_leads['Canal'] == canal]
        else: return {'Lead': 0, 'Workable Lead': 0}
    if 'Data de criação' not in df_c.columns: return {'Lead': 0, 'Workable Lead': 0}
    mask = (df_c['Data de criação'] >= start) & (df_c['Data de criação'] <= end)
    df_p = df_c[mask]; total = len(df_p)
    wl_col = next((c for c in df_c.columns if 'workable' in c.lower()), None)
    workable = 0
    if wl_col and len(df_p) > 0:
        try:
            workable = int(df_p[wl_col].apply(lambda v: 1 if str(v).strip().upper() in ('1','TRUE','VERDADEIRO') else 0).sum())
        except (ValueError, TypeError):
            workable = 0
    return {'Lead': total, 'Workable Lead': workable}


def calcular_metricas_dash_df(df_opps, df_leads, canal, re_start, re_end, mp_start, mp_end, sw_start, sw_end):
    df_c = filter_opps_df(df_opps, canal); fases = FASES_DASH[canal]
    tem_leads = any(f in fases for f in ('Lead', 'Workable Lead'))
    resultado = {}
    for periodo, start, end in [('realizado',re_start,re_end),('mes_passado',mp_start,mp_end),('semana',sw_start,sw_end)]:
        counts = {}
        if tem_leads:
            lc = count_leads_df(df_leads, canal, start, end)
            if 'Lead' in fases: counts['Lead'] = lc['Lead']
            if 'Workable Lead' in fases: counts['Workable Lead'] = lc['Workable Lead']
        for fase in fases:
            if fase in ('Lead', 'Workable Lead'): continue
            counts[fase] = count_stage_df(df_c, fase, start, end)
        resultado[periodo] = counts
    return resultado


def perc_mtd_ref(ref_date):
    sem = (ref_date.day - 1) // 7 + 1
    return DIST_MTD.get(ref_date.month, {}).get(sem, 1.0)


def badge_semaforo(conv_real, conv_plan=None):
    if conv_plan is None: return '#f3f4f6', '#4b5563'
    if conv_plan and conv_plan > 0:
        ratio = conv_real / conv_plan
        if ratio >= 1.00: return '#d1fae5', '#065f46'
        if ratio >= 0.95: return '#fef3c7', '#92400e'
        return '#fee2e2', '#991b1b'
    if conv_real >= 75: return '#d1fae5', '#065f46'
    if conv_real >= 50: return '#fef3c7', '#92400e'
    return '#fee2e2', '#991b1b'


def cor_numero_vs_meta(valor, meta):
    if not meta or meta == 0: return '#1a1a2e'
    r = valor / meta * 100
    if r >= 90: return '#065f46'
    if r >= 70: return '#92400e'
    return '#991b1b'


def kpi_cor(pct):
    if pct >= 20: return '#065f46'
    if pct >= 10: return '#92400e'
    return '#991b1b'


def fmt_num(v):
    if v is None or v == 0: return '—'
    if v >= 1000: return f"{int(v):,}".replace(',', '.')
    return str(int(v))

def gerar_dashboard_png(canal, metricas_plan, metricas_mes, metricas_sem, metricas_real,
                         ref=None, re_start=None, re_end=None, mp_start=None, mp_end=None, sw_start=None, sw_end=None,
                         pct_mtd_override=None):
    fases = FASES_DASH[canal]; convs = CONVERSOES_DASH[canal]; n_f = len(fases); n_c = len(convs)
    RH = 0.52; HDR_H = RH*1.30; SUB_H = RH*0.60; SEP_H = RH*0.40; CONV_H = RH*0.90
    fig_w = 10.0; fig_h = HDR_H + SUB_H + n_f*RH + SEP_H + n_c*CONV_H + 0.20
    fig, ax = plt.subplots(figsize=(fig_w, fig_h)); fig.patch.set_facecolor('white')
    ax.set_xlim(0, fig_w); ax.set_ylim(0, fig_h); ax.axis('off')
    cx = [0.0, 2.85, 4.72, 6.58, 8.38, fig_w]
    def mid(i): return (cx[i] + cx[i+1]) / 2
    HDR_TXT = ['#9ca3af','#0369a1','#1d4ed8','#6d28d9','#065f46']
    HDR_BG = ['white','#f0f9ff','#eff6ff','#faf5ff','#f0fdf4']
    CEL_BG = ['white','#f8fcff','#f8fbff','#fbf8ff','#f8fffe']
    HDR_LABELS = ['','PLANEJADO','MÊS PASSADO','SEMANA PASSADA','REALIZADO']
    _ref = ref or pd.Timestamp.now(); _re_start = re_start or pd.Timestamp.now().replace(day=1)
    _re_end = re_end or pd.Timestamp.now(); _mp_start = mp_start or pd.Timestamp.now().replace(day=1)
    _mp_end = mp_end or pd.Timestamp.now(); _sw_start = sw_start or pd.Timestamp.now().replace(day=1)
    _sw_end = sw_end or pd.Timestamp.now(); _pct_mtd = pct_mtd_override if pct_mtd_override is not None else perc_mtd_ref(_ref)
    HDR_DATES = ['', f"{MESES_PT[_ref.month]} {_ref.year} · MTD {_pct_mtd*100:.0f}%",
                 f"{_mp_start.strftime('%d/%m')} → {_mp_end.strftime('%d/%m')}",
                 f"{_sw_start.strftime('%d/%m')} → {_sw_end.strftime('%d/%m')}",
                 f"{_re_start.strftime('%d/%m')} → {_re_end.strftime('%d/%m')}"]
    y = fig_h
    y -= HDR_H
    for i in range(5):
        ax.add_patch(plt.Rectangle((cx[i],y),cx[i+1]-cx[i],HDR_H,facecolor=HDR_BG[i],edgecolor='none',zorder=1))
        ax.text(mid(i),y+HDR_H/2,HDR_LABELS[i],ha='center',va='center',fontsize=12,fontweight='bold',color=HDR_TXT[i],zorder=3)
    y -= SUB_H
    for i in range(5):
        ax.add_patch(plt.Rectangle((cx[i],y),cx[i+1]-cx[i],SUB_H,facecolor=HDR_BG[i],edgecolor='none',zorder=1))
        ax.text(mid(i),y+SUB_H/2,HDR_DATES[i],ha='center',va='center',fontsize=9,color='#9ca3af',zorder=3)
    ax.axhline(y,color='#e8eaf0',linewidth=1.2,zorder=2)
    all_metricas = [metricas_plan, metricas_mes, metricas_sem, metricas_real]
    def conv_plan_fase(ri):
        if ri == 0: return None
        p = metricas_plan.get(fases[ri]); pp = metricas_plan.get(fases[ri-1])
        return p/pp*100 if pp and pp > 0 and p else None
    for ri, fase in enumerate(fases):
        y -= RH; row_bg = '#f9fafb' if ri%2==1 else 'white'
        ax.add_patch(plt.Rectangle((cx[0],y),cx[1]-cx[0],RH,facecolor=row_bg,edgecolor='none',zorder=1))
        for ci in range(1,5):
            bg = CEL_BG[ci] if row_bg=='white' else '#f5f8fb' if ci<=2 else '#f7f4fc' if ci==3 else '#f4fcf7'
            ax.add_patch(plt.Rectangle((cx[ci],y),cx[ci+1]-cx[ci],RH,facecolor=bg,edgecolor='none',zorder=1))
        dot_c = DASH_DOT_CORES.get(fase,'#888')
        ax.add_patch(plt.Circle((0.22,y+RH/2),0.090,color=dot_c,zorder=3))
        is_nc = fase=="Novos contratos"
        ax.text(0.38,y+RH/2,fase,ha='left',va='center',fontsize=11,fontweight='bold' if is_nc else 'normal',color='#1a1a2e',zorder=3)
        pct_plan = conv_plan_fase(ri)
        for ci, metricas in enumerate(all_metricas):
            col_i = ci+1; val = metricas.get(fase); meta = metricas_plan.get(fase); is_plan_col = (col_i==1)
            pct = None
            if ri > 0:
                prev_val = metricas.get(fases[ri-1])
                if prev_val and prev_val > 0 and val: pct = val/prev_val*100
            if not val:
                ax.text(mid(col_i),y+RH/2,'—',ha='center',va='center',fontsize=10,color='#9ca3af',zorder=3); continue
            val_str = fmt_num(val); fw = 'bold'
            num_c = ('#0369a1' if is_nc else '#1a1a2e') if is_plan_col else cor_numero_vs_meta(val, meta)
            if pct is not None:
                col_w = cx[col_i+1]-cx[col_i]; val_x = cx[col_i]+col_w*0.36
                bw, bh = 0.46, RH*0.52; bx = val_x+0.07; by = y+(RH-bh)/2
                ax.text(val_x,y+RH/2,val_str,ha='right',va='center',fontsize=11,fontweight=fw,color=num_c,zorder=3)
                bg_b, fg_b = badge_semaforo(pct, None if is_plan_col else pct_plan)
                ax.add_patch(mpatches.FancyBboxPatch((bx,by),bw,bh,boxstyle="round,pad=0.03",facecolor=bg_b,edgecolor='none',zorder=3))
                ax.text(bx+bw/2,y+RH/2,f"{pct:.0f}%",ha='center',va='center',fontsize=9,fontweight='bold',color=fg_b,zorder=4)
            else:
                ax.text(mid(col_i),y+RH/2,val_str,ha='center',va='center',fontsize=11,fontweight=fw,color=num_c,zorder=3)
        ax.axhline(y,color='#f0f2f5',linewidth=0.5,zorder=2)
    y -= SEP_H; ax.axhline(y+SEP_H*0.5,color='#e8eaf0',linewidth=1.0,zorder=2)
    def calc_kpi(metricas, fase_num, fase_den):
        n = metricas.get(fase_num); d = metricas.get(fase_den)
        return d/n*100 if n and n > 0 and d else None
    for label, fase_num, fase_den in convs:
        y -= CONV_H
        ax.add_patch(plt.Rectangle((0,y),fig_w,CONV_H,facecolor='#f9fafb',edgecolor='none',zorder=1))
        ax.text(0.38,y+CONV_H/2,label,ha='left',va='center',fontsize=9.5,style='italic',color='#9ca3af',zorder=3)
        for ci, metricas in enumerate(all_metricas):
            col_i = ci+1; v = calc_kpi(metricas, fase_num, fase_den)
            if v is None: txt_c = '#9ca3af'; txt = '—'
            else: txt = f"{v:.1f}%"; txt_c = '#4b5563' if col_i==1 else kpi_cor(v)
            ax.text(mid(col_i),y+CONV_H/2,txt,ha='center',va='center',fontsize=9.5,style='italic',fontweight='bold',color=txt_c,zorder=3)
    for i in range(1,5): ax.axvline(cx[i],color='#e8eaf0',linewidth=0.6,zorder=2)
    ax.add_patch(mpatches.FancyBboxPatch((0.03,0.03),fig_w-0.06,fig_h-0.06,boxstyle="round,pad=0.03",facecolor='none',edgecolor='#e8eaf0',linewidth=1.5,zorder=5))
    plt.subplots_adjust(left=0,right=1,top=1,bottom=0)
    buf = io.BytesIO(); plt.savefig(buf,format='png',dpi=300,bbox_inches='tight',facecolor='white'); plt.close(); buf.seek(0); return buf.read()


# ── Slide manipulation ──

def remover_funis_existentes(slide):
    shapes_rem = [s for s in slide.shapes if s.shape_type == 13 and s.width > 1_500_000]
    for s in shapes_rem:
        sp = s._element; sp.getparent().remove(sp)

def add_img(slide, blob, pos):
    left, top, w, h = pos
    slide.shapes.add_picture(io.BytesIO(blob), Inches(left), Inches(top), Inches(w), Inches(h))

def fix_dates(slide, subs):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in subs.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)

# ══════════════════════════════════════════════════════════════
# HELPERS DE DATA
# ══════════════════════════════════════════════════════════════

def sexta_mais_recente(ref=None):
    d = ref or date.today()
    dias = (d.weekday() - 4) % 7
    return d - timedelta(days=dias)

def calcular_datas_auto():
    atual = sexta_mais_recente()
    sem_pass = atual - timedelta(days=7)
    mes_ant_mes = atual.month - 1 if atual.month > 1 else 12
    mes_ant_ano = atual.year if atual.month > 1 else atual.year - 1
    ultimo_dia = calendar.monthrange(mes_ant_ano, mes_ant_mes)[1]
    dia_ref = min(atual.day, ultimo_dia)
    mes_ant = date(mes_ant_ano, mes_ant_mes, dia_ref)
    return atual, sem_pass, mes_ant

# ══════════════════════════════════════════════════════════════
# TAXA DE JUROS — PROCESSAMENTO E RENDERIZAÇÃO
# ══════════════════════════════════════════════════════════════

TAXA_CANAIS_MAP = {
    'B2C': 'B2C',
    'Correspondente': 'Parceiros\nCorrespondentes',
    'Parceiro': 'Grandes\nParcerias',
    'Relacionamento': 'Relacionamento',
}
TAXA_CANAIS_ORDEM = ['B2C', 'Correspondente', 'Parceiro', 'Relacionamento']

METAS_TAXA_DEFAULT = {
    'B2C': 0.0138,
    'Correspondente': 0.0140,
    'Parceiro': 0.0140,
    'Relacionamento': 0.0140,
    'Geral': 0.0138,
}


def processar_taxas(file_bytes):
    """Processa base de taxas e retorna dict com resultados por canal."""
    import numpy as np
    df = pd.read_excel(io.BytesIO(file_bytes))

    # Passo 0a: Remover Derivadas
    mask_deriv = df['Nome da oportunidade'].str.contains('Derivada', case=False, na=False)
    df = df[~mask_deriv].copy()

    # Passo 0b: Remover Pré-Fixados (se coluna Modalidade existir)
    if 'Modalidade' in df.columns:
        mask_prefixado = df['Modalidade'].str.contains('Pré-Fixado|Pre-Fixado|PRÉ-FIXADO|PRE-FIXADO', case=False, na=False)
        df = df[~mask_prefixado].copy()

    # Passo 1: Chave de agrupamento por pessoa.
    # Usa CPF como chave primária — é imune a variações de acento, maiúscula,
    # espaçamento ou nome abreviado que o Salesforce introduz entre lançamentos
    # da mesma pessoa (ex.: "Jose Charlis Dias Sobrinho" vs "José Charlis dias
    # sobrinho", ou "GABRIEL VELOSO CAMBRAIA" vs "Gabriel Cambraia" — mesmo CPF
    # nos dois casos). Antes, o agrupamento por nome digitado deixava essas
    # linhas soltas, fora da consolidação de renegociação, distorcendo a média.
    # Quando o CPF vier vazio, cai para o nome normalizado (sem acento) como
    # reserva, para não descartar a linha.
    import unicodedata
    def _normalizar_nome(x):
        if pd.isna(x):
            return ''
        nome = str(x).split(' - ')[0].strip().upper()
        nome = unicodedata.normalize('NFKD', nome).encode('ASCII', 'ignore').decode('ASCII')
        return nome

    def _chave_pessoa(row):
        cpf = str(row.get('CPF', '')).strip()
        if cpf and cpf.lower() != 'nan':
            return f"CPF:{cpf}"
        return f"NOME:{_normalizar_nome(row['Nome da oportunidade'])}"

    df['chave_pessoa'] = df.apply(_chave_pessoa, axis=1)

    # Passo 2 & 3: Agrupar e consolidar renegociações
    def consolidar_grupo(group):
        if len(group) <= 1:
            return group
        reneg_mask = group['Nome da oportunidade'].str.contains('Renegociação', case=False, na=False)
        if not reneg_mask.any():
            return group
        reneg_idx = group[reneg_mask].index[0]
        # Soma ignorando NaN
        soma_valor = group['Valor do Derivado'].fillna(0).sum()
        result = group.copy()
        result.loc[reneg_idx, 'Valor do Derivado'] = soma_valor
        outras = result.index != reneg_idx
        result.loc[outras, 'Valor do Derivado'] = 0
        result.loc[outras, 'Taxa'] = 0
        result.loc[outras, 'Taxa bonificada'] = np.nan
        return result

    df = df.groupby('chave_pessoa', group_keys=False).apply(consolidar_grupo)

    # Passo 4: Calcular Taxa Nova
    def calc_taxa_nova(row):
        if pd.isna(row['Valor do Derivado']) or row['Valor do Derivado'] <= 0:
            return 0
        taxa = row['Taxa'] if pd.notna(row['Taxa']) else 0
        taxa_bon = row['Taxa bonificada']
        if pd.notna(taxa_bon):
            return taxa_bon * 0.79 + taxa * 0.21
        else:
            return taxa

    df['Taxa Nova'] = df.apply(calc_taxa_nova, axis=1)

    # Filtrar valor > 0
    df_pos = df[df['Valor do Derivado'] > 0].copy()

    # Passo 5, 6, 7: Calcular por canal
    resultados = {}
    for canal in TAXA_CANAIS_ORDEM:
        dc = df_pos[df_pos['Canal'] == canal]
        if len(dc) == 0:
            resultados[canal] = {'contratual': 0, 'ponderada': 0}
            continue
        sum_val = dc['Valor do Derivado'].sum()
        contratual = (dc['Taxa'] * dc['Valor do Derivado']).sum() / sum_val
        ponderada = (dc['Taxa Nova'] * dc['Valor do Derivado']).sum() / sum_val
        resultados[canal] = {'contratual': contratual, 'ponderada': ponderada}

    # Geral
    sum_val_total = df_pos['Valor do Derivado'].sum()
    if sum_val_total > 0:
        contratual_geral = (df_pos['Taxa'] * df_pos['Valor do Derivado']).sum() / sum_val_total
        ponderada_geral = (df_pos['Taxa Nova'] * df_pos['Valor do Derivado']).sum() / sum_val_total
    else:
        contratual_geral = ponderada_geral = 0
    resultados['Geral'] = {'contratual': contratual_geral, 'ponderada': ponderada_geral}

    return resultados


def _cor_semaforo_taxa(valor, meta):
    """Para taxas: verde se >= meta (bom, banco cobra mais), vermelho se < meta (ruim)."""
    if valor >= meta:
        return '#16A34A'
    return '#DC2626'


def _cor_semaforo_ticket(pct_realizado):
    """Para ticket: verde 90-110%, amarelo >110%, vermelho <90%."""
    if pct_realizado >= 1.10:
        return '#EAB308'  # amarelo — muito acima
    if pct_realizado >= 0.90:
        return '#16A34A'  # verde — na meta
    return '#DC2626'      # vermelho — abaixo


def _gerar_tabela_bari(titulo, subtitulo, headers, canais_data, mes_nome, fmt_func):
    """
    Gera imagem PNG de tabela estilo Bari (fundo branco, bordas pretas, espaçado).
    Sem título/logo — já está na apresentação.
    canais_data: list of (label, [(valor, cor_bg), ...]) — uma tupla por coluna de dados
    fmt_func: função que formata o valor para string
    """
    n_cols = len(headers)
    n_rows = len(canais_data)
    fig_w, fig_h = 9.6, 0.60 + n_rows * 0.78
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    bg_color = 'white'
    fig.patch.set_facecolor(bg_color)
    ax.set_xlim(0, fig_w); ax.set_ylim(0, fig_h); ax.axis('off')
    ax.set_facecolor(bg_color)

    # Layout de colunas
    label_w = 2.10
    gap = 0.12
    total_data_w = fig_w - label_w - 0.50
    data_w = (total_data_w - gap * (n_cols - 2)) / (n_cols - 1)
    col_starts = [0.30]
    for i in range(1, n_cols):
        col_starts.append(label_w + 0.20 + (i - 1) * (data_w + gap))

    row_h = 0.66
    row_gap = 0.10
    header_y = fig_h - 0.30
    cell_pad = 0.06

    # Cabeçalhos
    for i, txt in enumerate(headers):
        if i == 0: continue
        ax.text(col_starts[i] + data_w / 2, header_y, txt,
                ha='center', va='center', fontsize=11, fontweight='bold', color='#6b7280', zorder=3)

    # Linhas de dados
    for ri, (label, cells) in enumerate(canais_data):
        y = header_y - (ri + 1) * (row_h + row_gap) + row_gap * 0.3

        # Separador antes de Geral
        if label == 'Geral':
            ax.axhline(y + row_h + row_gap * 0.4, xmin=0.03, xmax=0.97,
                       color='#d1d5db', linewidth=0.8, zorder=2)

        # Label
        ax.text(col_starts[0] + label_w - 0.15, y + row_h / 2, label,
                ha='right', va='center', fontsize=11.5,
                fontweight='bold' if label == 'Geral' else 'normal',
                color='#1a1a2e', zorder=3, linespacing=1.3)

        # Cells
        for ci, (valor, cor_bg) in enumerate(cells):
            cx = col_starts[ci + 1]
            txt = fmt_func(valor)
            # Fundo com borda preta
            ax.add_patch(mpatches.FancyBboxPatch(
                (cx + cell_pad, y + cell_pad), data_w - 2 * cell_pad, row_h - 2 * cell_pad,
                boxstyle="round,pad=0.05", facecolor=cor_bg,
                edgecolor='#1a1a2e', linewidth=1.5, zorder=3))
            # Cor do texto
            txt_color = '#2D3142' if cor_bg in ('#EEF0F4', '#f3f4f6', '#e5e7eb', '#E0E7EF') else 'white'
            ax.text(cx + data_w / 2, y + row_h / 2, txt,
                    ha='center', va='center', fontsize=14, fontweight='bold',
                    color=txt_color, family='monospace', zorder=4)

    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor=bg_color)
    plt.close()
    buf.seek(0)
    return buf.read()


def gerar_taxa_png(resultados, metas_taxa, mes_nome):
    """Gera imagem PNG da tabela de taxa de juros."""
    labels_map = {
        'B2C': 'B2C', 'Correspondente': 'Parceiros\nCorrespondentes',
        'Parceiro': 'Grandes\nParcerias', 'Relacionamento': 'Relacionamento', 'Geral': 'Geral',
    }
    canais_data = []
    for canal in TAXA_CANAIS_ORDEM + ['Geral']:
        res = resultados.get(canal, {'contratual': 0, 'ponderada': 0})
        meta = metas_taxa.get(canal, 0.0138)
        cells = [
            (meta, '#EEF0F4'),  # Meta — fundo cinza claro
            (res['contratual'], _cor_semaforo_taxa(res['contratual'], meta)),
            (res['ponderada'], _cor_semaforo_taxa(res['ponderada'], meta)),
        ]
        canais_data.append((labels_map[canal], cells))

    def fmt_taxa(v):
        return f"{v*100:.2f}%".replace('.', ',')

    return _gerar_tabela_bari(
        titulo=f'Taxa de juros Varejo  -  {mes_nome}',
        subtitulo='Apenas novos contratos ou aditivos (sem derivadas)',
        headers=['', 'Meta', 'Contratual', 'Ponderada'],
        canais_data=canais_data,
        mes_nome=mes_nome,
        fmt_func=fmt_taxa,
    )


def processar_ticket_medio(taxas_bytes, contratos_bytes):
    """Calcula ticket médio por canal. Retorna dict {canal: ticket}."""
    df_taxa = pd.read_excel(io.BytesIO(taxas_bytes))
    df_contr = pd.read_excel(io.BytesIO(contratos_bytes))

    # Volume por canal (COM derivadas)
    vol_por_canal = df_taxa.groupby('Canal')['Valor do Derivado'].sum()
    contr_por_canal = df_contr.groupby('Canal').size()

    resultados = {}
    total_vol = 0
    total_contr = 0
    for canal in TAXA_CANAIS_ORDEM:
        vol = vol_por_canal.get(canal, 0)
        contr = contr_por_canal.get(canal, 0)
        ticket = vol / contr if contr > 0 else 0
        resultados[canal] = ticket
        total_vol += vol
        total_contr += contr

    resultados['Geral'] = total_vol / total_contr if total_contr > 0 else 0
    return resultados


def gerar_ticket_png(resultados_ticket, metas_ticket, mes_nome):
    """Gera imagem PNG da tabela de ticket médio."""
    labels_map = {
        'B2C': 'B2C', 'Correspondente': 'Parceiros\nCorrespondentes',
        'Parceiro': 'Grandes\nParcerias', 'Relacionamento': 'Relacionamento', 'Geral': 'Geral',
    }
    canais_data = []
    for canal in TAXA_CANAIS_ORDEM + ['Geral']:
        ticket = resultados_ticket.get(canal, 0)
        meta = metas_ticket.get(canal, 250000)
        pct = ticket / meta if meta > 0 else 0
        cells = [
            (meta, '#EEF0F4'),      # Meta — cinza
            (ticket, '#EEF0F4'),    # Realizado — cinza igual meta
            (pct, _cor_semaforo_ticket(pct)),  # % realizado
        ]
        canais_data.append((labels_map[canal], cells))

    def fmt_ticket(v):
        if isinstance(v, float) and v < 10:
            # É percentual
            return f"{v*100:.0f}%"
        return f"{v:,.0f}".replace(',', '.')

    return _gerar_tabela_bari(
        titulo=f'Análise de ticket médio novos contratos - {mes_nome}',
        subtitulo='',
        headers=['', 'Meta', 'Realizado', '% realizado'],
        canais_data=canais_data,
        mes_nome=mes_nome,
        fmt_func=fmt_ticket,
    )

METAS_ORIG_MENSAL = {
    #        B2C         Correspondente(PC)  Parceiro(GP)    Relacionamento  Carteira
    1:  {'B2C': 8000000,  'Correspondente': 15400000, 'Parceiro': 17000000, 'Relacionamento': 5800000, 'Compra de carteira': 3000000},
    2:  {'B2C': 9000000,  'Correspondente': 16500000, 'Parceiro': 17000000, 'Relacionamento': 6000000, 'Compra de carteira': 3000000},
    3:  {'B2C': 10800000, 'Correspondente': 20000000, 'Parceiro': 19800000, 'Relacionamento': 7000000, 'Compra de carteira': 3000000},
    4:  {'B2C': 9200000,  'Correspondente': 17300000, 'Parceiro': 20800000, 'Relacionamento': 4800000, 'Compra de carteira': 3000000},
    5:  {'B2C': 9600000,  'Correspondente': 17900000, 'Parceiro': 21200000, 'Relacionamento': 4800000, 'Compra de carteira': 3000000},
    6:  {'B2C': 10200000, 'Correspondente': 19000000, 'Parceiro': 22800000, 'Relacionamento': 5200000, 'Compra de carteira': 3000000},
    7:  {'B2C': 11800000, 'Correspondente': 21600000, 'Parceiro': 25400000, 'Relacionamento': 6100000, 'Compra de carteira': 3000000},
    8:  {'B2C': 11300000, 'Correspondente': 21300000, 'Parceiro': 24600000, 'Relacionamento': 5800000, 'Compra de carteira': 3000000},
    9:  {'B2C': 11000000, 'Correspondente': 20500000, 'Parceiro': 24000000, 'Relacionamento': 5600000, 'Compra de carteira': 3000000},
    10: {'B2C': 11700000, 'Correspondente': 21700000, 'Parceiro': 26300000, 'Relacionamento': 5900000, 'Compra de carteira': 3000000},
    11: {'B2C': 11400000, 'Correspondente': 20800000, 'Parceiro': 25400000, 'Relacionamento': 5400000, 'Compra de carteira': 3000000},
    12: {'B2C': 12400000, 'Correspondente': 22500000, 'Parceiro': 27500000, 'Relacionamento': 5800000, 'Compra de carteira': 3000000},
}

METAS_CONTR_MENSAL = {
    #        B2C  Correspondente(PC)  Parceiro(GP)  Relacionamento
    1:  {'B2C': 37, 'Correspondente': 33, 'Parceiro': 76, 'Relacionamento': 36},
    2:  {'B2C': 44, 'Correspondente': 35, 'Parceiro': 76, 'Relacionamento': 33},
    3:  {'B2C': 48, 'Correspondente': 42, 'Parceiro': 86, 'Relacionamento': 43},
    4:  {'B2C': 40, 'Correspondente': 36, 'Parceiro': 89, 'Relacionamento': 29},
    5:  {'B2C': 41, 'Correspondente': 37, 'Parceiro': 89, 'Relacionamento': 29},
    6:  {'B2C': 43, 'Correspondente': 39, 'Parceiro': 94, 'Relacionamento': 31},
    7:  {'B2C': 49, 'Correspondente': 44, 'Parceiro': 103, 'Relacionamento': 36},
    8:  {'B2C': 46, 'Correspondente': 43, 'Parceiro': 98, 'Relacionamento': 34},
    9:  {'B2C': 44, 'Correspondente': 41, 'Parceiro': 94, 'Relacionamento': 33},
    10: {'B2C': 46, 'Correspondente': 43, 'Parceiro': 101, 'Relacionamento': 34},
    11: {'B2C': 44, 'Correspondente': 41, 'Parceiro': 96, 'Relacionamento': 31},
    12: {'B2C': 47, 'Correspondente': 44, 'Parceiro': 102, 'Relacionamento': 33},
}

# Defaults para fallback (usa abril)
METAS_ORIG_DEFAULT = METAS_ORIG_MENSAL[4]
METAS_CONTR_DEFAULT = METAS_CONTR_MENSAL[4]


def processar_originacao(taxas_bytes, contratos_bytes, valor_carteira_manual=0):
    """Calcula valor originado (COM derivadas) e novos contratos por canal.

    valor_carteira_manual: valor de Compra de carteira digitado no site.
    Como esse canal NÃO vem da base do Salesforce, o valor é informado
    manualmente toda semana. Quando > 0, é usado na linha Compra de carteira;
    quando 0, a linha aparece zerada (vira '-' no slide).
    Compra de carteira não conta como novo contrato → contratos sempre 0.
    """
    df_taxa = pd.read_excel(io.BytesIO(taxas_bytes))
    df_contr = pd.read_excel(io.BytesIO(contratos_bytes))

    vol_por_canal = df_taxa.groupby('Canal')['Valor do Derivado'].sum()
    contr_por_canal = df_contr.groupby('Canal').size()

    resultados = {}
    for canal in TAXA_CANAIS_ORDEM:
        resultados[canal] = {
            'valor': vol_por_canal.get(canal, 0),
            'contratos': contr_por_canal.get(canal, 0),
        }

    # Compra de carteira: valor manual do site tem prioridade.
    # Se não informado (0), tenta a base (que normalmente não tem o canal → 0).
    val_carteira = valor_carteira_manual if valor_carteira_manual and valor_carteira_manual > 0 \
        else vol_por_canal.get('Compra de carteira', 0)
    resultados['Compra de carteira'] = {'valor': val_carteira, 'contratos': 0}

    return resultados


def gerar_originacao_png(resultados, metas_orig, metas_contr):
    """Gera imagem PNG da tabela de Originação + Novos Contratos — design premium para diretoria."""
    labels_map = {
        'B2C': 'B2C', 'Correspondente': 'Parceiros\nCorrespondentes',
        'Parceiro': 'Grandes\nParcerias', 'Relacionamento': 'Relacionamento',
    }

    canais = TAXA_CANAIS_ORDEM

    # Calcular totais
    varejo_val = sum(resultados.get(c, {}).get('valor', 0) for c in canais)
    varejo_contr = sum(resultados.get(c, {}).get('contratos', 0) for c in canais)
    varejo_meta_val = sum(metas_orig.get(c, 0) for c in canais)
    varejo_meta_contr = sum(metas_contr.get(c, 0) for c in canais)

    carteira = resultados.get('Compra de carteira', {})
    carteira_val = carteira.get('valor', 0)
    carteira_meta = metas_orig.get('Compra de carteira', 3000000)

    total_val = varejo_val + carteira_val
    total_meta_val = varejo_meta_val + carteira_meta

    # Layout
    fig_w = 10.0
    row_h = 0.58
    row_gap = 0.08
    sep_h = 0.18
    n_data_rows = 7  # 4 canais + varejo + carteira + total
    n_seps = 2
    header_h = 0.80
    fig_h = header_h + n_data_rows * (row_h + row_gap) + n_seps * sep_h + 0.20

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor('white')
    ax.set_xlim(0, fig_w)
    ax.set_ylim(0, fig_h)
    ax.axis('off')

    # Column layout
    label_w = 1.35
    section_gap = 0.30
    orig_total_w = 4.95
    contr_total_w = 3.30
    orig_x = label_w + 0.10
    contr_x = orig_x + orig_total_w + section_gap

    # Proportions within each section: meta | realizado | %
    orig_props = [0.38, 0.38, 0.24]
    contr_props = [0.30, 0.30, 0.40]

    def fmt_val(v):
        if v == 0: return '-'
        return f"{v:,.0f}".replace(',', '.')

    def fmt_pct(real, meta):
        if meta == 0 or real == 0: return '-'
        return f"{real/meta*100:.0f}%"

    def draw_bar(x, y, total_w, h, values, props, bold=False, highlight=False, hide=False):
        """Desenha uma barra conectada (Meta + Realizado + %) com bordas compartilhadas."""
        if hide:
            return
        bg = '#e8ecf2' if highlight else '#f1f5f9'
        ec = '#d0d5dd'
        fs = 15 if not bold else 16
        fw = 'bold' if bold else 'normal'

        # Barra principal (retângulo com borda arredondada)
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, y), total_w, h,
            boxstyle="round,pad=0.03",
            facecolor=bg, edgecolor=ec, linewidth=1.0, zorder=3))

        # Divisórias verticais internas
        cx = x
        for i, (val, prop) in enumerate(zip(values, props)):
            cell_w = total_w * prop
            # Divisória (exceto antes da primeira)
            if i > 0:
                ax.plot([cx, cx], [y + 0.04, y + h - 0.04],
                        color=ec, linewidth=0.8, zorder=4)

            # Texto
            tc = '#0f172a' if val != '-' else '#cbd5e1'
            ax.text(cx + cell_w / 2, y + h / 2, str(val),
                    ha='center', va='center', fontsize=fs, fontweight=fw,
                    color=tc, zorder=5)
            cx += cell_w

    # Section headers
    y = fig_h - 0.25
    ax.text(orig_x + orig_total_w / 2, y, 'ORIGINAÇÃO',
            ha='center', va='center', fontsize=15, fontweight='bold',
            color='#0f172a', zorder=3)
    ax.text(contr_x + contr_total_w / 2, y, 'NOVOS CONTRATOS',
            ha='center', va='center', fontsize=15, fontweight='bold',
            color='#0f172a', zorder=3)

    # Sub-headers
    y -= 0.35
    for i, txt in enumerate(['Meta', 'Realizado', '']):
        ox = orig_x + sum(orig_props[:i]) * orig_total_w + orig_props[i] * orig_total_w / 2
        ax.text(ox, y, txt, ha='center', va='center',
                fontsize=9.5, fontweight='bold', color='#94a3b8', zorder=3)
        cx_ = contr_x + sum(contr_props[:i]) * contr_total_w + contr_props[i] * contr_total_w / 2
        ax.text(cx_, y, txt, ha='center', va='center',
                fontsize=9.5, fontweight='bold', color='#94a3b8', zorder=3)

    # Data rows
    y -= 0.25
    row_data = []
    for canal in canais:
        res = resultados.get(canal, {'valor': 0, 'contratos': 0})
        meta_v = metas_orig.get(canal, 0)
        meta_c = metas_contr.get(canal, 0)
        row_data.append((labels_map[canal], meta_v, res['valor'], meta_c, res['contratos'], False, False, True))

    row_data.append(('SEP', 0, 0, 0, 0, False, False, False))
    row_data.append(('Varejo', varejo_meta_val, varejo_val, varejo_meta_contr, varejo_contr, True, True, True))
    row_data.append(('Compra de\ncarteira', carteira_meta, carteira_val, 0, 0, False, False, False))
    row_data.append(('SEP', 0, 0, 0, 0, False, False, False))
    row_data.append(('TOTAL', total_meta_val, total_val, 0, 0, True, True, False))

    for label, meta_v, real_v, meta_c, real_c, bold, highlight, show_contr in row_data:
        if label == 'SEP':
            ax.axhline(y, xmin=0.01, xmax=0.99, color='#0f172a', linewidth=1.8, zorder=2)
            y -= sep_h
            continue

        # Label
        ax.text(label_w - 0.05, y - row_h / 2, label,
                ha='right', va='center', fontsize=11.5,
                fontweight='bold' if bold else 'normal',
                color='#0f172a', zorder=3, linespacing=1.2)

        # Originação bar
        ry = y - row_h
        pct_v = fmt_pct(real_v, meta_v)
        draw_bar(orig_x, ry, orig_total_w, row_h,
                 [fmt_val(meta_v), fmt_val(real_v) if real_v > 0 else '-', pct_v],
                 orig_props, bold=bold, highlight=highlight)

        # Novos Contratos bar
        if show_contr:
            pct_c = fmt_pct(real_c, meta_c)
            draw_bar(contr_x, ry, contr_total_w, row_h,
                     [fmt_val(meta_c), fmt_val(real_c) if real_c > 0 else '-', pct_c],
                     contr_props, bold=bold, highlight=highlight)

        y -= (row_h + row_gap)

    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════
# SNAPSHOT — SALVAR E CARREGAR RETRATOS HISTÓRICOS NO GITHUB
# ══════════════════════════════════════════════════════════════

GITHUB_REPO = "CharlesFaria/Apresenta-o-Semanal-Diretoria-"
SNAPSHOTS_DIR = "snapshots"


def _get_github_token():
    """Obtém o token do GitHub dos secrets do Streamlit."""
    try:
        return st.secrets.get("GITHUB_TOKEN", None)
    except Exception:
        return None


def _github_headers():
    token = _get_github_token()
    if not token:
        return None
    return {
        "Authorization": f"token {token}",
        "Accept": "application/vnd.github.v3+json",
    }


def salvar_snapshot_github(data_ref, snapshot_data):
    """
    Salva um snapshot no GitHub como JSON.
    Arquivo: snapshots/2026-04-25.json
    """
    headers = _github_headers()
    if not headers:
        return False, "Token GitHub não configurado"

    filename = f"{SNAPSHOTS_DIR}/{data_ref.strftime('%Y-%m-%d')}.json"
    url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{filename}"

    content_b64 = base64.b64encode(json.dumps(snapshot_data, ensure_ascii=False, default=str).encode()).decode()

    # Verificar se já existe (precisa do SHA para atualizar)
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code == 200:
            sha = resp.json().get("sha")
            body = {
                "message": f"Atualizar snapshot {data_ref.strftime('%d/%m/%Y')}",
                "content": content_b64,
                "sha": sha,
            }
        else:
            body = {
                "message": f"Snapshot {data_ref.strftime('%d/%m/%Y')}",
                "content": content_b64,
            }

        resp = requests.put(url, headers=headers, json=body, timeout=15)
        if resp.status_code in (200, 201):
            return True, f"Snapshot salvo: {filename}"
        else:
            return False, f"Erro ao salvar: {resp.status_code} - {resp.text[:200]}"
    except Exception as e:
        return False, f"Erro: {str(e)}"


def listar_snapshots_github():
    """Lista todos os snapshots disponíveis no GitHub."""
    headers = _github_headers()
    if not headers:
        return []

    url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{SNAPSHOTS_DIR}"
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code != 200:
            return []
        files = resp.json()
        datas = []
        for f in files:
            nome = f.get("name", "")
            if nome.endswith(".json"):
                try:
                    d = datetime.strptime(nome.replace(".json", ""), "%Y-%m-%d").date()
                    datas.append(d)
                except ValueError:
                    pass
        return sorted(datas)
    except Exception:
        return []


def carregar_snapshot_github(data_ref):
    """Carrega um snapshot específico do GitHub."""
    headers = _github_headers()
    if not headers:
        return None

    filename = f"{SNAPSHOTS_DIR}/{data_ref.strftime('%Y-%m-%d')}.json"
    url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{filename}"
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code != 200:
            return None
        content_b64 = resp.json().get("content", "")
        content = base64.b64decode(content_b64).decode()
        return json.loads(content)
    except Exception:
        return None


def buscar_snapshot_mais_proximo(data_alvo, tolerancia_dias=3):
    """
    Busca o snapshot mais próximo da data_alvo (±3 dias = mesma semana).
    Retorna (data_snapshot, dados) ou (None, None) se não encontrar.
    """
    datas = listar_snapshots_github()
    if not datas:
        return None, None

    melhor = None
    menor_diff = timedelta(days=999)
    for d in datas:
        diff = abs(d - data_alvo)
        if diff < menor_diff and diff <= timedelta(days=tolerancia_dias):
            menor_diff = diff
            melhor = d

    if melhor is None:
        return None, None

    dados = carregar_snapshot_github(melhor)
    return melhor, dados


def montar_snapshot(data_ref, cache_funis, fases_ativas, resultados_taxa=None,
                     resultados_ticket=None, resultados_orig=None):
    """
    Monta o dict do snapshot com todos os dados da apresentação.
    cache_funis: dict de (canal, tipo, "atual") → bytes PNG (base64 encoded)
    """
    snapshot = {
        "data": data_ref.isoformat(),
        "gerado_em": datetime.now().isoformat(),
        "funis": {},
        "taxa": resultados_taxa,
        "ticket": resultados_ticket,
        "originacao": resultados_orig,
    }

    # Salvar funis como base64
    for (canal, tipo, periodo), png_bytes in cache_funis.items():
        if periodo == "atual":
            key = f"{canal}_{tipo}"
            snapshot["funis"][key] = base64.b64encode(png_bytes).decode()

    return snapshot


def extrair_funil_do_snapshot(snapshot, canal, tipo):
    """Extrai bytes PNG de um funil do snapshot."""
    key = f"{canal}_{tipo}"
    b64 = snapshot.get("funis", {}).get(key)
    if b64:
        return base64.b64decode(b64)
    return None


# ══════════════════════════════════════════════════════════════
# PERSISTÊNCIA DE METAS NO GITHUB
# ══════════════════════════════════════════════════════════════

METAS_FILE = "metas.json"


@st.cache_data(ttl=300)  # Cache por 5 minutos
def carregar_metas_github():
    """Carrega metas salvas do GitHub."""
    headers = _github_headers()
    if not headers:
        return {}
    url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{METAS_FILE}"
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code != 200:
            return {}
        content_b64 = resp.json().get("content", "")
        content = base64.b64decode(content_b64).decode()
        return json.loads(content)
    except Exception:
        return {}


def salvar_metas_github(metas_data):
    """Salva metas no GitHub como JSON. Retorna (ok: bool, mensagem: str)."""
    headers = _github_headers()
    if not headers:
        return False, "Token GitHub não configurado"
    url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{METAS_FILE}"
    content_b64 = base64.b64encode(json.dumps(metas_data, ensure_ascii=False).encode()).decode()
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code == 200:
            sha = resp.json().get("sha")
            body = {"message": "Atualizar metas", "content": content_b64, "sha": sha}
        elif resp.status_code == 404:
            body = {"message": "Criar metas", "content": content_b64}
        else:
            return False, f"Erro ao verificar arquivo: {resp.status_code} - {resp.text[:200]}"

        resp = requests.put(url, headers=headers, json=body, timeout=15)
        if resp.status_code in (200, 201):
            carregar_metas_github.clear()
            return True, "Metas salvas com sucesso"
        else:
            return False, f"Erro ao salvar: {resp.status_code} - {resp.text[:200]}"
    except Exception as e:
        return False, f"Erro de conexão: {str(e)}"


# ══════════════════════════════════════════════════════════════
# PROCESSAMENTO PRINCIPAL
# ══════════════════════════════════════════════════════════════

def processar_tudo(pptx_bytes, base_funil_bytes, base_dash_bytes, base_leads_bytes,
                   plan_bytes, data_atual, data_sem_pass, data_mes_ant, progress_bar, status_text,
                   dist_mtd_user=None, semana_atual=3, taxas_bytes=None, metas_taxa=None,
                   contratos_bytes=None, metas_ticket=None, metas_orig=None, metas_contr=None,
                   valor_carteira=0):
    """Processa tudo e retorna (bytes_pptx, lista_de_logs).

    valor_carteira: valor de Compra de carteira digitado manualmente no site
    (nasce 0 a cada execução; não persiste e não vem da base do Salesforce).
    """
    # Atualiza DIST_MTD com valores do usuário
    if dist_mtd_user:
        for m in range(1, 13):
            DIST_MTD[m] = dist_mtd_user

    logs = []
    def log(msg):
        logs.append(msg)

    total_steps = 50
    step = 0
    def advance(n=1, msg=None):
        nonlocal step
        step += n
        progress_bar.progress(min(step / total_steps, 1.0))
        if msg:
            status_text.markdown(f"<span style='font-size:13px;color:#64748b'>{msg}</span>", unsafe_allow_html=True)

    log("╔══════════════════════════════════════════════╗")
    log("║   Automatização Slides Semanais — Bari       ║")
    log("╚══════════════════════════════════════════════╝")
    log("")

    STR_ATUAL = data_atual.strftime("%d/%m"); MES_ATUAL = MESES_PT[data_atual.month]
    STR_SEM_PASS = data_sem_pass.strftime("%d/%m"); MES_SEM_PASS = MESES_PT[data_sem_pass.month]
    STR_MES_ANT = data_mes_ant.strftime("%d/%m"); MES_ANT = MESES_PT[data_mes_ant.month]
    log(f"📅 Datas: atual={STR_ATUAL} | sem_pass={STR_SEM_PASS} | mes_ant={STR_MES_ANT}")

    # Planejamento
    METAS_2026 = {}
    if plan_bytes:
        advance(2, "📅 Carregando planejamento...")
        METAS_2026 = carregar_planejamento(plan_bytes)
        total_p = sum(len(f) for cd in METAS_2026.values() for f in cd.values())
        log(f"\n📅 Planejamento: {total_p} valores carregados")

    # Base funil
    advance(3, "📊 Carregando base do funil...")
    rows = carregar_base(base_funil_bytes)
    log(f"\n📊 Base do funil: {len(rows)} oportunidades")

    # ── Buscar snapshot do mês anterior ──
    snapshot_mes_ant = None
    snapshot_data_real = None
    advance(1, "🔍 Buscando snapshot do mês anterior...")
    snap_date, snap_data = buscar_snapshot_mais_proximo(data_mes_ant)
    if snap_data:
        snapshot_mes_ant = snap_data
        snapshot_data_real = snap_date
        log(f"\n📸 Snapshot encontrado: {snap_date.strftime('%d/%m/%Y')} (pedido: {data_mes_ant.strftime('%d/%m/%Y')})")
    else:
        log(f"\n📸 Nenhum snapshot próximo de {data_mes_ant.strftime('%d/%m/%Y')} — calculando da base")

    # Gerar funis
    advance(2, "🎨 Gerando funis...")
    log("\n🎨 Gerando funis...")
    cache = {}
    combos = list(dict.fromkeys((c, t) for _, c, t, _ in SLIDES_FUNIL))
    periodos = [("atual", data_atual), ("sem_pass", data_sem_pass), ("mes_ant", data_mes_ant)]
    # Coletar fases ativas por canal (para legendas dinâmicas)
    fases_ativas = {}
    for canal, tipo in combos:
        for chave, data_ref in periodos:
            # Se é mes_ant E temos snapshot, usar a imagem salva
            if chave == "mes_ant" and snapshot_mes_ant:
                funil_salvo = extrair_funil_do_snapshot(snapshot_mes_ant, canal, tipo)
                if funil_salvo:
                    cache[(canal, tipo, chave)] = funil_salvo
                    continue
            # Caso contrário, calcular normalmente
            count, volume = retrato_funil(rows, canal, data_ref)
            mes = MESES_PT[data_ref.month]
            data_label = f"{mes} ({data_ref.strftime('%d/%m')})"
            cache[(canal, tipo, chave)] = gerar_funil_png(count, volume, canal, tipo, data_label)
            # Coleta fases com valor > 0 para a legenda
            for f in FASES_NOMES:
                if count.get(f, 0) > 0 or volume.get(f, 0) > 0:
                    fases_ativas.setdefault(canal, set()).add(f)
        c_atual, v_atual = retrato_funil(rows, canal, data_atual)
        if tipo == "volume":
            lbl = f"R${sum(v_atual.values())/1e6:.1f}M"
        else:
            lbl = f"{int(sum(c_atual.values()))} propostas"
        log(f"  ✅ {canal:12s} | {tipo:10s} | atual: {lbl}")
        advance(2)

    # Gerar legendas dinâmicas por canal (mantém a ordem original, só exclui zerados)
    legendas = {}
    for canal in set(c for c, _ in combos):
        fases_canal = [f for f in FASES_NOMES if f in fases_ativas.get(canal, set())]
        legendas[canal] = gerar_legenda_png(fases_canal if fases_canal else FASES_NOMES)
    log("  ✅ Legendas")
    advance(1)

    # Apresentação
    advance(2, "📂 Carregando apresentação...")
    prs = Presentation(io.BytesIO(pptx_bytes))
    log(f"\n📂 Apresentação: {len(prs.slides)} slides")

    # Atualizar slides funil
    advance(2, "🔄 Atualizando slides dos funis...")
    log("\n🔄 Atualizando slides dos funis...")
    nomes = {"B2C": "B2C", "GP": "Grandes Parcerias", "PC": "Correspondentes (PC)", "Comercial": "Comercial s/ Rel"}
    for num, canal, tipo, comp in SLIDES_FUNIL:
        idx = num - 1
        if idx >= len(prs.slides):
            log(f"  ⚠️  Slide {num} não existe"); continue
        slide = prs.slides[idx]; nome = nomes.get(canal, canal)
        periodo_esq = "mes_ant" if comp == "mensal" else "sem_pass"
        img_esq = cache[(canal, tipo, periodo_esq)]; img_dir = cache[(canal, tipo, "atual")]
        leg = legendas.get(canal, legendas.get("B2C", gerar_legenda_png()))
        remover_funis_existentes(slide)
        add_img(slide, img_esq, POS_ESQ); add_img(slide, img_dir, POS_DIR); add_img(slide, leg, POS_LEGENDA)
        if comp == "mensal":
            fix_dates(slide, {STR_MES_ANT: STR_MES_ANT, MES_ANT: MES_ANT, STR_SEM_PASS: STR_ATUAL, MES_SEM_PASS: MES_ATUAL})
        else:
            fix_dates(slide, {STR_SEM_PASS: STR_SEM_PASS, STR_ATUAL: STR_ATUAL})
        log(f"  Slide {num:2d} — {nome:20s} | {tipo:10s} | {comp}")
        advance(1)

    # Dashboards
    if base_dash_bytes:
        advance(2, "📋 Gerando dashboards...")
        log("\n📋 Gerando dashboards...")
        df_opps = carregar_opps_df(base_dash_bytes)
        log(f"   {len(df_opps)} oportunidades (dash)")
        df_leads = carregar_leads_df(base_leads_bytes)
        if df_leads is not None:
            log(f"   {len(df_leads)} leads")
        else:
            log("   ⚠️ Sem base de leads — Lead/WL zerados")
        REF, re_start, re_end, mp_start, mp_end, sw_start, sw_end = calcular_ref_periodos(df_opps)
        log(f"   REF={REF.strftime('%d/%m/%Y')} | realizado={re_start.strftime('%d/%m')}→{re_end.strftime('%d/%m')}")
        for num, canal in SLIDES_DASH:
            idx = num - 1
            if idx >= len(prs.slides):
                log(f"  ⚠️  Slide {num} não existe"); continue
            mes = REF.month; pct_mtd = DIST_MTD.get(mes, {}).get(semana_atual, perc_mtd_ref(REF))
            plan_raw = METAS_2026.get(canal, {}).get(mes, {})
            metricas_plan = {f: (round(plan_raw[f]*pct_mtd) if plan_raw.get(f) else None) for f in FASES_DASH[canal]}
            resultado = calcular_metricas_dash_df(df_opps, df_leads, canal, re_start, re_end, mp_start, mp_end, sw_start, sw_end)
            png = gerar_dashboard_png(canal, metricas_plan, resultado['mes_passado'], resultado['semana'], resultado['realizado'],
                                       ref=REF, re_start=re_start, re_end=re_end, mp_start=mp_start, mp_end=mp_end, sw_start=sw_start, sw_end=sw_end,
                                       pct_mtd_override=pct_mtd)
            slide = prs.slides[idx]; remover_funis_existentes(slide); add_img(slide, png, POS_DASH)
            log(f"  Slide {num:2d} — {canal} ✅")
            advance(2)
    else:
        log("\n⚠️ Sem base do dashboard — pulando dashboards")

    # ── Slide de Originação + Novos Contratos (slide 5) ──
    if taxas_bytes and contratos_bytes:
        advance(2, "📊 Gerando slide de originação...")
        log("\n📊 Gerando slide de Originação + Novos Contratos...")
        try:
            resultados_orig = processar_originacao(taxas_bytes, contratos_bytes, valor_carteira_manual=valor_carteira)
            _metas_o = metas_orig or METAS_ORIG_DEFAULT
            _metas_c = metas_contr or METAS_CONTR_DEFAULT
            png_orig = gerar_originacao_png(resultados_orig, _metas_o, _metas_c)

            slide_orig_idx = 4  # slide 5
            if slide_orig_idx < len(prs.slides):
                slide = prs.slides[slide_orig_idx]
                remover_funis_existentes(slide)
                add_img(slide, png_orig, (0.20, 1.02, 9.35, 4.38))
                log(f"  Slide 5 — Originação + Novos Contratos ✅")
            else:
                log(f"  ⚠️ Slide 5 não existe")

            cart = resultados_orig.get('Compra de carteira', {})
            if cart.get('valor', 0) > 0:
                log(f"     {'Compra de carteira':20s} | Valor: R${cart.get('valor',0):,.0f} (manual)")
            for canal in TAXA_CANAIS_ORDEM:
                res = resultados_orig.get(canal, {})
                log(f"     {canal:20s} | Valor: R${res.get('valor',0):,.0f} | Contratos: {res.get('contratos',0)}")
        except Exception as e:
            log(f"  ❌ Erro ao processar originação: {str(e)}")
    else:
        if not taxas_bytes or not contratos_bytes:
            log("\n⚠️ Sem base de taxas ou contratos — pulando slide de Originação")

    # ── Slide de Taxa de Juros ──
    if taxas_bytes:
        advance(2, "💰 Gerando slide de taxas...")
        log("\n💰 Gerando slide de Taxa de Juros...")
        try:
            resultados_taxa = processar_taxas(taxas_bytes)
            _metas = metas_taxa or METAS_TAXA_DEFAULT
            mes_nome = MESES_PT[data_atual.month]
            png_taxa = gerar_taxa_png(resultados_taxa, _metas, mes_nome)

            # Slide 7 (índice 6)
            slide_taxa_idx = 6  # slide 7
            if slide_taxa_idx < len(prs.slides):
                slide = prs.slides[slide_taxa_idx]
                remover_funis_existentes(slide)
                add_img(slide, png_taxa, (0.63, 1.32, 7.73, 4.19))
                log(f"  Slide 7 — Taxa de Juros ✅")
            else:
                log(f"  ⚠️ Slide 7 não existe")

            # Log dos valores
            for canal in TAXA_CANAIS_ORDEM + ['Geral']:
                res = resultados_taxa.get(canal, {})
                c = res.get('contratual', 0) * 100
                p = res.get('ponderada', 0) * 100
                log(f"     {canal:20s} | Contratual: {c:.2f}% | Ponderada: {p:.2f}%")
        except Exception as e:
            log(f"  ❌ Erro ao processar taxas: {str(e)}")
    else:
        log("\n⚠️ Sem base de taxas — pulando slide de Taxa de Juros")

    # ── Slide de Ticket Médio (slide 8) ──
    if taxas_bytes and contratos_bytes:
        advance(2, "📋 Gerando slide de ticket médio...")
        log("\n📋 Gerando slide de Ticket Médio...")
        try:
            resultados_ticket = processar_ticket_medio(taxas_bytes, contratos_bytes)
            _metas_tk = metas_ticket or {'B2C': 220882, 'Correspondente': 471614, 'Parceiro': 226959, 'Relacionamento': 162619, 'Geral': 270518}
            mes_nome = MESES_PT[data_atual.month]
            png_ticket = gerar_ticket_png(resultados_ticket, _metas_tk, mes_nome)

            slide_ticket_idx = 7  # slide 8
            if slide_ticket_idx < len(prs.slides):
                slide = prs.slides[slide_ticket_idx]
                remover_funis_existentes(slide)
                add_img(slide, png_ticket, (0.63, 1.32, 7.73, 4.19))
                log(f"  Slide 8 — Ticket Médio ✅")
            else:
                log(f"  ⚠️ Slide 8 não existe")

            for canal in TAXA_CANAIS_ORDEM + ['Geral']:
                ticket = resultados_ticket.get(canal, 0)
                log(f"     {canal:20s} | Ticket: R${ticket:,.0f}")
        except Exception as e:
            log(f"  ❌ Erro ao processar ticket médio: {str(e)}")
    else:
        if not taxas_bytes or not contratos_bytes:
            log("\n⚠️ Sem base de taxas ou contratos — pulando slide de Ticket Médio")

    # Salvar
    advance(2, "💾 Salvando apresentação...")
    output = io.BytesIO(); prs.save(output); output.seek(0)

    # ── Salvar snapshot automaticamente ──
    advance(1, "📸 Salvando snapshot...")
    try:
        # Coletar resultados para o snapshot
        snap_taxas = None
        snap_ticket = None
        snap_orig = None
        try:
            if taxas_bytes:
                snap_taxas = processar_taxas(taxas_bytes)
                # Converter valores float para serializable
                snap_taxas = {k: {kk: float(vv) for kk, vv in v.items()} for k, v in snap_taxas.items()}
        except: pass
        try:
            if taxas_bytes and contratos_bytes:
                snap_ticket = processar_ticket_medio(taxas_bytes, contratos_bytes)
                snap_ticket = {k: float(v) for k, v in snap_ticket.items()}
        except: pass
        try:
            if taxas_bytes and contratos_bytes:
                snap_orig = processar_originacao(taxas_bytes, contratos_bytes, valor_carteira_manual=valor_carteira)
                snap_orig = {k: {kk: float(vv) for kk, vv in v.items()} for k, v in snap_orig.items()}
        except: pass

        snapshot = montar_snapshot(
            data_ref=data_atual,
            cache_funis=cache,
            fases_ativas=fases_ativas,
            resultados_taxa=snap_taxas,
            resultados_ticket=snap_ticket,
            resultados_orig=snap_orig,
        )
        ok, msg = salvar_snapshot_github(data_atual, snapshot)
        if ok:
            log(f"\n📸 {msg}")
        else:
            log(f"\n⚠️ Snapshot não salvo: {msg}")
    except Exception as e:
        log(f"\n⚠️ Erro ao salvar snapshot: {str(e)}")

    log("\n✅ Apresentação gerada com sucesso!")
    log("🎉 Pronto! Baixe o arquivo abaixo.")
    progress_bar.progress(1.0)
    status_text.markdown("<span style='font-size:13px;color:#2563EB;font-weight:700'>✅ Concluído!</span>", unsafe_allow_html=True)

    return output.getvalue(), logs


# ══════════════════════════════════════════════════════════════
# INTERFACE
# ══════════════════════════════════════════════════════════════

def render_file_card(title, subtitle, icon, accent, file_obj, required=False):
    """Renders info about a file card (purely visual, uploader is separate)."""
    has_file = file_obj is not None
    badge = ""
    if required:
        badge = f'<span class="badge-{"ok" if has_file else "req"}">{"OK" if has_file else "OBRIGATÓRIO"}</span>'

    icon_class = f"{accent}" if has_file else "empty"
    file_info = ""
    if has_file and file_obj is not None:
        size = f"({file_obj.size / 1024:.0f} KB)" if hasattr(file_obj, 'size') else ""
        name = file_obj.name if hasattr(file_obj, 'name') else "arquivo"
        file_info = f'<span class="upload-filename">{name}</span> <span style="color:#94a3b8;font-size:11px">{size}</span>'
    else:
        file_info = subtitle

    st.markdown(f"""
    <div class="upload-card {'upload-card-'+accent+' has-file file-'+accent if has_file else ''}">
        <div style="display:flex;align-items:center;gap:12px">
            <div class="upload-icon {icon_class}" style="color:white;font-size:16px">{icon if not has_file else '✓'}</div>
            <div style="flex:1">
                <div style="display:flex;align-items:center;gap:8px">
                    <span class="upload-title">{title}</span>
                    {badge}
                </div>
                <div class="upload-sub">{file_info}</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)


def main():
    # ── Header ──
    st.markdown("""
    <div class="bari-header">
        <div class="bari-header-inner">
            <div>
                <div class="bari-logo">bari<span>.</span></div>
                <div class="bari-header-title">Slides Semanais — Diretoria Comercial</div>
            </div>
            <div class="bari-header-right">
                <div class="bari-header-dot"></div>
                <span>Gerador de Apresentações</span>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Step 1: Bases ──
    st.markdown("""
    <div class="step-header">
        <div class="step-num">1</div>
        <span class="step-title">Carregar bases</span>
        <span class="step-sub">Exporte do Salesforce e arraste aqui</span>
    </div>
    """, unsafe_allow_html=True)

    # Links do Salesforce
    with st.expander("🔗 Links dos relatórios no Salesforce"):
        st.markdown("""
        Clique no link, exporte como **.xlsx** e suba no campo correspondente abaixo.

        | Base | Link |
        |------|------|
        | 📈 Oportunidades | [Abrir no Salesforce](https://barigui.lightning.force.com/lightning/r/Report/00OTT000007PrY92AK/view?queryScope=userFolders) |
        | 👥 Leads | [Abrir no Salesforce](https://barigui.lightning.force.com/lightning/r/Report/00OTT0000060P6D2AU/view?queryScope=userFolders) |
        | 💰 Taxas | [Abrir no Salesforce](https://barigui.lightning.force.com/lightning/r/Report/00OHY000000JqQi2AK/view) |
        | 📋 Contratos | [Abrir no Salesforce](https://barigui.lightning.force.com/lightning/r/Report/00OHY000000JrVZ2A0/view) |
        """)

    # Verificar se tem apresentação modelo no repositório
    MODELO_PATH = os.path.join(os.path.dirname(__file__), "apresentacao_modelo.pptx")
    tem_modelo = os.path.exists(MODELO_PATH)

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("**📈 Base de Oportunidades** `OBRIGATÓRIO`", help="Exportação completa do Salesforce")
        st.caption("Atualizar Entrada nas Fases.xlsx — use a base completa (período amplo)")
        f_opps = st.file_uploader("Base Oportunidades", type=["xlsx"], key="f_opps", label_visibility="collapsed")
    with col2:
        if tem_modelo:
            st.markdown("**📑 Apresentação Modelo** ✅ pré-carregada")
            st.caption("Usando apresentação do repositório. Suba outra se quiser substituir.")
            f_pptx_upload = st.file_uploader("Apresentação", type=["pptx"], key="f_pptx", label_visibility="collapsed")
            # Se o usuário subiu outra, usa ela; senão usa a do repositório
            if f_pptx_upload:
                f_pptx = f_pptx_upload
            else:
                f_pptx = MODELO_PATH  # será tratado como path no processar_tudo
        else:
            st.markdown("**📑 Apresentação Modelo** `OBRIGATÓRIO`", help="Arquivo .pptx da diretoria")
            st.caption("Arquivo .pptx atual da diretoria")
            f_pptx = st.file_uploader("Apresentação", type=["pptx"], key="f_pptx", label_visibility="collapsed")

    col3, col4 = st.columns(2)
    with col3:
        st.markdown("**👥 Base de Leads**")
        st.caption("Entradas nas Fases Leads.xlsx — para Lead/Workable Lead no dashboard")
        f_leads = st.file_uploader("Base Leads", type=["xlsx"], key="f_leads", label_visibility="collapsed")
    with col4:
        st.markdown("**💰 Base de Taxas**")
        st.caption("Valor efetivado com % de meta — para slide de Taxa de Juros")
        f_taxas = st.file_uploader("Base Taxas", type=["xlsx"], key="f_taxas", label_visibility="collapsed")

    col5, col6 = st.columns(2)
    with col5:
        st.markdown("**📋 Base de Contratos**")
        st.caption("Novos contratos com % de meta — para slide de Ticket Médio")
        f_contratos = st.file_uploader("Base Contratos", type=["xlsx"], key="f_contratos", label_visibility="collapsed")
    with col6:
        with st.expander("🎯 Planejamento — muda pouco"):
            st.caption("Planejamento.xlsx — metas mensais por canal")
            f_plan = st.file_uploader("Planejamento", type=["xlsx"], key="f_plan", label_visibility="collapsed")

    st.markdown('<div class="soft-divider"></div>', unsafe_allow_html=True)

    # ── Step 2: Datas ──
    st.markdown("""
    <div class="step-header">
        <div class="step-num">2</div>
        <span class="step-title">Conferir datas</span>
    </div>
    """, unsafe_allow_html=True)

    auto_atual, auto_sem, auto_mes = calcular_datas_auto()
    usar_auto = st.toggle("Calcular automaticamente", value=True, key="auto_datas")

    if usar_auto:
        data_atual, data_sem_pass, data_mes_ant = auto_atual, auto_sem, auto_mes
        c1, c2, c3 = st.columns(3)
        with c1:
            st.markdown(f"""<div class="date-card date-card-green date-green">
                <div class="date-label">DATA ATUAL (SEXTA)</div>
                <div class="date-value">{data_atual.strftime('%d/%m/%Y')}</div>
            </div>""", unsafe_allow_html=True)
        with c2:
            st.markdown(f"""<div class="date-card date-card-blue date-blue">
                <div class="date-label">SEMANA PASSADA</div>
                <div class="date-value">{data_sem_pass.strftime('%d/%m/%Y')}</div>
            </div>""", unsafe_allow_html=True)
        with c3:
            st.markdown(f"""<div class="date-card date-card-purple date-purple">
                <div class="date-label">MÊS ANTERIOR</div>
                <div class="date-value">{data_mes_ant.strftime('%d/%m/%Y')}</div>
            </div>""", unsafe_allow_html=True)
    else:
        c1, c2, c3 = st.columns(3)
        with c1:
            data_atual = st.date_input("Data atual (sexta)", value=auto_atual, key="d_atual")
        with c2:
            data_sem_pass = st.date_input("Semana passada", value=auto_sem, key="d_sem")
        with c3:
            data_mes_ant = st.date_input("Mês anterior", value=auto_mes, key="d_mes")

    st.markdown('<div class="soft-divider"></div>', unsafe_allow_html=True)

    # ── Distribuição MTD ──
    with st.expander("📊 Distribuição semanal do planejamento (% por semana · dias úteis)"):
        st.caption("Digite o % de cada semana (baseado em dias úteis) e selecione em qual semana estamos.")

        semana_atual = st.radio(
            "Estamos na semana:",
            options=[1, 2, 3, 4, 5],
            format_func=lambda x: f"Sem {x}",
            horizontal=True,
            index=2,
            key="semana_atual"
        )

        mc1, mc2, mc3, mc4, mc5 = st.columns(5)
        with mc1:
            pct_s1 = st.number_input("Sem 1 (%)", min_value=0.0, max_value=100.0, value=7.2, step=0.1, format="%.1f", key="mtd1")
        with mc2:
            pct_s2 = st.number_input("Sem 2 (%)", min_value=0.0, max_value=100.0, value=21.0, step=0.1, format="%.1f", key="mtd2")
        with mc3:
            pct_s3 = st.number_input("Sem 3 (%)", min_value=0.0, max_value=100.0, value=28.8, step=0.1, format="%.1f", key="mtd3")
        with mc4:
            pct_s4 = st.number_input("Sem 4 (%)", min_value=0.0, max_value=100.0, value=21.0, step=0.1, format="%.1f", key="mtd4")
        with mc5:
            pct_s5 = st.number_input("Sem 5 (%)", min_value=0.0, max_value=100.0, value=22.0, step=0.1, format="%.1f", key="mtd5")

        # Calcular acumulado
        pcts = [pct_s1, pct_s2, pct_s3, pct_s4, pct_s5]
        soma = sum(pcts)
        acumulados = []
        ac = 0
        for p in pcts:
            ac += p
            acumulados.append(ac)

        # Mostrar resumo com semana atual destacada
        acum_parts = []
        for i, a in enumerate(acumulados):
            if i + 1 == semana_atual:
                acum_parts.append(f"**[{a:.1f}%]**")
            else:
                acum_parts.append(f"{a:.1f}%")
        acum_str = " → ".join(acum_parts)

        if abs(soma - 100.0) < 0.5:
            st.success(f"✅ Soma: {soma:.1f}%  ·  Acumulado: {acum_str}  ·  MTD atual (sem {semana_atual}): **{acumulados[semana_atual-1]:.1f}%**")
        else:
            st.warning(f"⚠️ Soma: {soma:.1f}% (deveria ser 100%)  ·  Acumulado: {acum_str}")

    # Converter para dict acumulado normalizado (0 a 1)
    ac = 0
    dist_mtd_user = {}
    for i, p in enumerate(pcts, 1):
        ac += p
        dist_mtd_user[i] = ac / 100

    st.markdown('<div class="soft-divider"></div>', unsafe_allow_html=True)

    # ── Carregar metas salvas do GitHub ──
    metas_salvas = carregar_metas_github()

    # ── Metas de Taxa ──
    _tx_saved = metas_salvas.get('taxa', {})
    metas_taxa = dict(METAS_TAXA_DEFAULT)
    if f_taxas:
        with st.expander("💰 Metas de Taxa de Juros"):
            st.caption("Ajuste as metas de taxa por canal (formato decimal, ex: 1.38 = 1,38%)")
            mt1, mt2, mt3, mt4, mt5 = st.columns(5)
            with mt1:
                metas_taxa['B2C'] = st.number_input("B2C (%)", min_value=0.0, max_value=5.0, value=float(_tx_saved.get('B2C', 1.38)), step=0.01, format="%.2f", key="meta_b2c") / 100
            with mt2:
                metas_taxa['Correspondente'] = st.number_input("Corresp. (%)", min_value=0.0, max_value=5.0, value=float(_tx_saved.get('Correspondente', 1.40)), step=0.01, format="%.2f", key="meta_corresp") / 100
            with mt3:
                metas_taxa['Parceiro'] = st.number_input("Parceiro (%)", min_value=0.0, max_value=5.0, value=float(_tx_saved.get('Parceiro', 1.40)), step=0.01, format="%.2f", key="meta_parc") / 100
            with mt4:
                metas_taxa['Relacionamento'] = st.number_input("Relac. (%)", min_value=0.0, max_value=5.0, value=float(_tx_saved.get('Relacionamento', 1.40)), step=0.01, format="%.2f", key="meta_rel") / 100
            with mt5:
                metas_taxa['Geral'] = st.number_input("Geral (%)", min_value=0.0, max_value=5.0, value=float(_tx_saved.get('Geral', 1.38)), step=0.01, format="%.2f", key="meta_geral") / 100

    # ── Metas de Ticket Médio (por mês) ──
    METAS_TICKET_MENSAL = {
        1:  {'B2C': 216216, 'Correspondente': 466667, 'Parceiro': 223684, 'Relacionamento': 161111, 'Geral': 253846},
        2:  {'B2C': 204545, 'Correspondente': 471429, 'Parceiro': 223684, 'Relacionamento': 181818, 'Geral': 257979},
        3:  {'B2C': 225000, 'Correspondente': 476190, 'Parceiro': 230233, 'Relacionamento': 162971, 'Geral': 263014},
        4:  {'B2C': 230000, 'Correspondente': 480556, 'Parceiro': 233708, 'Relacionamento': 165517, 'Geral': 268557},
        5:  {'B2C': 234146, 'Correspondente': 483784, 'Parceiro': 238202, 'Relacionamento': 165517, 'Geral': 272959},
        6:  {'B2C': 237209, 'Correspondente': 487179, 'Parceiro': 242553, 'Relacionamento': 167742, 'Geral': 276329},
        7:  {'B2C': 240816, 'Correspondente': 490909, 'Parceiro': 246602, 'Relacionamento': 167742, 'Geral': 279741},
        8:  {'B2C': 245652, 'Correspondente': 495349, 'Parceiro': 251020, 'Relacionamento': 169444, 'Geral': 285068},
        9:  {'B2C': 250000, 'Correspondente': 500000, 'Parceiro': 255319, 'Relacionamento': 170588, 'Geral': 288208},
        10: {'B2C': 254348, 'Correspondente': 504651, 'Parceiro': 260396, 'Relacionamento': 169697, 'Geral': 292857},
        11: {'B2C': 259091, 'Correspondente': 507317, 'Parceiro': 264363, 'Relacionamento': 173529, 'Geral': 297170},
        12: {'B2C': 263830, 'Correspondente': 511364, 'Parceiro': 269608, 'Relacionamento': 175758, 'Geral': 301770},
    }
    mes_ref = data_atual.month if usar_auto else data_atual.month
    metas_ticket_default = METAS_TICKET_MENSAL.get(mes_ref, METAS_TICKET_MENSAL[4])
    _tk_saved = metas_salvas.get('ticket', {})
    metas_ticket = dict(metas_ticket_default)

    if f_contratos and f_taxas:
        with st.expander("📋 Metas de Ticket Médio"):
            st.caption(f"Metas para {MESES_PT[mes_ref]} (carregadas automaticamente — ajuste se necessário)")
            tk1, tk2, tk3, tk4, tk5 = st.columns(5)
            with tk1:
                metas_ticket['B2C'] = st.number_input("B2C (R$)", min_value=0, value=int(_tk_saved.get('B2C', metas_ticket_default['B2C'])), step=1000, key="meta_tk_b2c")
            with tk2:
                metas_ticket['Correspondente'] = st.number_input("PC (R$)", min_value=0, value=int(_tk_saved.get('Correspondente', metas_ticket_default['Correspondente'])), step=1000, key="meta_tk_corresp")
            with tk3:
                metas_ticket['Parceiro'] = st.number_input("GP (R$)", min_value=0, value=int(_tk_saved.get('Parceiro', metas_ticket_default['Parceiro'])), step=1000, key="meta_tk_parc")
            with tk4:
                metas_ticket['Relacionamento'] = st.number_input("Rel (R$)", min_value=0, value=int(_tk_saved.get('Relacionamento', metas_ticket_default['Relacionamento'])), step=1000, key="meta_tk_rel")
            with tk5:
                metas_ticket['Geral'] = st.number_input("Geral (R$)", min_value=0, value=int(_tk_saved.get('Geral', metas_ticket_default['Geral'])), step=1000, key="meta_tk_geral")

    # ── Metas de Originação (por mês) ──
    _ov_saved = metas_salvas.get('originacao', {})
    _oc_saved = metas_salvas.get('contratos', {})
    metas_orig_default = METAS_ORIG_MENSAL.get(mes_ref, METAS_ORIG_DEFAULT)
    metas_contr_default = METAS_CONTR_MENSAL.get(mes_ref, METAS_CONTR_DEFAULT)
    metas_orig = dict(metas_orig_default)
    metas_contr = dict(metas_contr_default)
    valor_carteira = 0  # default; sobrescrito pelo campo abaixo se o expander existir.
                        # Nasce 0 a cada execução: não persiste nem herda do mês anterior.
    if f_contratos and f_taxas:
        with st.expander("📊 Metas de Originação + Novos Contratos"):
            st.caption(f"Metas para {MESES_PT[mes_ref]} (carregadas automaticamente)")
            st.caption("Metas de valor originado (R$)")
            ov1, ov2, ov3, ov4, ov5 = st.columns(5)
            with ov1:
                metas_orig['B2C'] = st.number_input("B2C (R$)", min_value=0, value=int(_ov_saved.get('B2C', metas_orig_default['B2C'])), step=100000, key="meta_ov_b2c")
            with ov2:
                metas_orig['Correspondente'] = st.number_input("PC (R$)", min_value=0, value=int(_ov_saved.get('Correspondente', metas_orig_default['Correspondente'])), step=100000, key="meta_ov_corresp")
            with ov3:
                metas_orig['Parceiro'] = st.number_input("GP (R$)", min_value=0, value=int(_ov_saved.get('Parceiro', metas_orig_default['Parceiro'])), step=100000, key="meta_ov_parc")
            with ov4:
                metas_orig['Relacionamento'] = st.number_input("Rel (R$)", min_value=0, value=int(_ov_saved.get('Relacionamento', metas_orig_default['Relacionamento'])), step=100000, key="meta_ov_rel")
            with ov5:
                metas_orig['Compra de carteira'] = st.number_input("Carteira (R$)", min_value=0, value=int(_ov_saved.get('Compra de carteira', metas_orig_default.get('Compra de carteira', 3000000))), step=100000, key="meta_ov_cart")

            st.caption("Metas de novos contratos (quantidade)")
            oc1, oc2, oc3, oc4 = st.columns(4)
            with oc1:
                metas_contr['B2C'] = st.number_input("B2C (qtd)", min_value=0, value=int(_oc_saved.get('B2C', metas_contr_default['B2C'])), step=1, key="meta_oc_b2c")
            with oc2:
                metas_contr['Correspondente'] = st.number_input("PC (qtd)", min_value=0, value=int(_oc_saved.get('Correspondente', metas_contr_default['Correspondente'])), step=1, key="meta_oc_corresp")
            with oc3:
                metas_contr['Parceiro'] = st.number_input("GP (qtd)", min_value=0, value=int(_oc_saved.get('Parceiro', metas_contr_default['Parceiro'])), step=1, key="meta_oc_parc")
            with oc4:
                metas_contr['Relacionamento'] = st.number_input("Rel (qtd)", min_value=0, value=int(_oc_saved.get('Relacionamento', metas_contr_default['Relacionamento'])), step=1, key="meta_oc_rel")

            # ── Valor REALIZADO de Compra de carteira (manual, não é meta) ──
            # Como esse canal não vem do Salesforce, o valor é digitado aqui toda
            # semana. value=0 fixo: nasce zerado a cada execução, NÃO persiste e
            # NÃO herda do mês anterior. Nos meses sem compra de carteira, deixe 0
            # e a linha aparece como "-" no slide.
            st.markdown("---")
            st.caption("💼 Valor realizado de Compra de carteira (R$) — preencha só nos meses que tiveram")
            valor_carteira = st.number_input(
                "Compra de carteira — realizado (R$)",
                min_value=0,
                value=0,
                step=100000,
                key="valor_carteira_realizado",
                help="Compra de carteira não vem do Salesforce: digite o valor aqui. "
                     "Deixe 0 nos meses sem compra de carteira (a linha vira '-' no slide)."
            )

    # ── Botão Salvar Metas ──
    if st.button("💾 Salvar metas no servidor", use_container_width=True):
        todas_metas = {
            'taxa': {k: round(v * 100, 2) for k, v in metas_taxa.items()},
            'ticket': {k: int(v) for k, v in metas_ticket.items()},
            'originacao': {k: int(v) for k, v in metas_orig.items()},
            'contratos': {k: int(v) for k, v in metas_contr.items()},
        }
        ok, msg = salvar_metas_github(todas_metas)
        if ok:
            st.success("✅ Metas salvas! Agora persistem entre sessões.")
        else:
            st.warning(f"⚠️ Não foi possível salvar: {msg}")

    st.markdown('<div class="soft-divider"></div>', unsafe_allow_html=True)

    # ── Snapshots ──
    with st.expander("📸 Snapshots históricos (automático)"):
        token_ok = _get_github_token() is not None
        if token_ok:
            st.markdown("✅ **GitHub Token configurado** — snapshots serão salvos automaticamente")
            datas_snap = listar_snapshots_github()
            if datas_snap:
                st.caption(f"{len(datas_snap)} snapshot(s) salvos:")
                for d in datas_snap[-10:]:  # últimos 10
                    st.markdown(f"  · `{d.strftime('%d/%m/%Y')}`")
            else:
                st.caption("Nenhum snapshot salvo ainda. O primeiro será criado ao gerar a apresentação.")
        else:
            st.markdown("⚠️ **GitHub Token não configurado** — snapshots não serão salvos")
            st.markdown("""
            Para ativar os snapshots automáticos:
            1. No GitHub, vá em **Settings** → **Developer settings** → **Personal access tokens** → **Tokens (classic)**
            2. Gere um token com permissão **repo**
            3. No Streamlit Cloud, vá em **Manage app** → **Settings** → **Secrets**
            4. Adicione: `GITHUB_TOKEN = "seu_token_aqui"`
            """)

    st.markdown('<div class="soft-divider"></div>', unsafe_allow_html=True)

    # ── Step 3: Gerar ──
    can_generate = f_opps is not None and f_pptx is not None

    step_bg = "var(--bari-blue)" if can_generate else "var(--bari-gray-200)"
    step_color = "white" if can_generate else "var(--bari-gray-400)"
    title_color = "var(--bari-navy)" if can_generate else "var(--bari-gray-400)"

    st.markdown(f"""
    <div class="step-header">
        <div class="step-num" style="background:{step_bg};color:{step_color}">3</div>
        <span class="step-title" style="color:{title_color}">Gerar apresentação</span>
    </div>
    """, unsafe_allow_html=True)

    if can_generate:
        # Summary
        pptx_nome = f_pptx.name if hasattr(f_pptx, 'name') else 'apresentacao_modelo.pptx'
        parts = [f"Oportunidades ({f_opps.name})"]
        if f_leads: parts.append("Leads")
        if f_taxas: parts.append("Taxas")
        if f_contratos: parts.append("Contratos")
        if f_plan: parts.append("Planejamento")
        st.markdown(f"""<div class="summary-box">
            <strong style="color:#0A1628">Resumo:</strong> {' + '.join(parts)} → <strong style="color:#2563EB">{pptx_nome}</strong>
        </div>""", unsafe_allow_html=True)

    if not can_generate:
        st.markdown("""<div class="note-box">
            ⏳ Carregue pelo menos a <strong>Base de Oportunidades</strong> para continuar.
        </div>""", unsafe_allow_html=True)
        return

    if st.button("🚀  Gerar Apresentação", type="primary", use_container_width=True, disabled=not can_generate):
        progress_bar = st.progress(0)
        status_text = st.empty()

        try:
            # Ler pptx (pode ser upload ou path do repositório)
            if isinstance(f_pptx, str):
                with open(f_pptx, 'rb') as f:
                    pptx_bytes = f.read()
            else:
                pptx_bytes = f_pptx.read()

            opps_bytes = f_opps.read()
            taxas_bytes_read = f_taxas.read() if f_taxas else None
            contratos_bytes_read = f_contratos.read() if f_contratos else None
            result_bytes, log_lines = processar_tudo(
                pptx_bytes=pptx_bytes,
                base_funil_bytes=opps_bytes,
                base_dash_bytes=opps_bytes,
                base_leads_bytes=f_leads.read() if f_leads else None,
                plan_bytes=f_plan.read() if f_plan else None,
                data_atual=data_atual,
                data_sem_pass=data_sem_pass,
                data_mes_ant=data_mes_ant,
                progress_bar=progress_bar,
                status_text=status_text,
                dist_mtd_user=dist_mtd_user,
                semana_atual=semana_atual,
                taxas_bytes=taxas_bytes_read,
                metas_taxa=metas_taxa,
                contratos_bytes=contratos_bytes_read,
                metas_ticket=metas_ticket,
                metas_orig=metas_orig,
                metas_contr=metas_contr,
                valor_carteira=valor_carteira,
            )

            # Success
            nome_saida = f"Apresentacao_{MESES_PT[data_atual.month]}_{data_atual.strftime('%d%m%Y')}.pptx"

            st.markdown("""<div class="success-banner">
                <div style="font-size:40px;margin-bottom:8px">🎉</div>
                <div style="font-size:18px;font-weight:800;color:#2563EB;margin-bottom:4px">Apresentação pronta!</div>
            </div>""", unsafe_allow_html=True)

            st.download_button(
                label=f"📥  Baixar {nome_saida}",
                data=result_bytes,
                file_name=nome_saida,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True,
            )

            # Log (colapsável)
            with st.expander("📋 Ver log completo"):
                log_html = ""
                for l in log_lines:
                    css = "log-normal"
                    if "✅" in l: css = "log-success"
                    elif "⚠" in l: css = "log-warning"
                    elif "❌" in l: css = "log-error"
                    elif "🎉" in l or "📅" in l or "📊" in l or "📂" in l or "📋" in l or "🔄" in l or "💾" in l: css = "log-info"
                    log_html += f'<div class="{css}">{l}</div>'
                st.markdown(f'<div class="log-panel">{log_html}</div>', unsafe_allow_html=True)

        except Exception as e:
            st.error(f"Erro durante o processamento: {str(e)}")
            import traceback
            st.code(traceback.format_exc())

    # Nota sobre leads
    if can_generate and not f_leads and not st.session_state.get('_generated'):
        st.markdown("""<div class="note-box" style="margin-top:14px">
            <strong>Nota:</strong> Sem a base de Leads, os campos Lead e Workable Lead nos dashboards ficarão zerados.
        </div>""", unsafe_allow_html=True)


if __name__ == "__main__":
    main()
